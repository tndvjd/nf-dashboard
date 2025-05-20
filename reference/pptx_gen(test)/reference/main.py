import sys
import os
import shutil
import tempfile
import asyncio
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QGridLayout, 
                             QLabel, QLineEdit, QTextEdit, QPushButton, QComboBox, QMessageBox, 
                             QProgressBar, QFileDialog, QGroupBox, QStatusBar, QToolTip, QWhatsThis)
from PyQt6.QtCore import QObject, pyqtSignal, Qt, QTimer, QEvent
from PyQt6.QtGui import QFont, QPainter, QColor
from qasync import QEventLoop, asyncSlot
from pptx import Presentation
from scraper import scrape_broker_info, scrape_property_info, download_floor_plan, download_complex_photo
from ppt_utils import add_slide_note_with_broker_info, update_pptx_template, update_cover_slide, add_images_to_rectangles, add_new_slide
from gen_map import get_coordinates, get_static_map_url, download_image
from translator import DeepLTranslator

import logging
logging.getLogger('qasync').setLevel(logging.WARNING)
logging.getLogger('asyncio').setLevel(logging.WARNING)

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

class ProgressWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        layout = QVBoxLayout(self)
        self.progress_bar = QProgressBar()
        self.status_label = QLabel("대기 중...")
        self.step_label = QLabel("단계: -")
        layout.addWidget(self.progress_bar)
        layout.addWidget(self.status_label)
        layout.addWidget(self.step_label)

    def update_progress(self, value, status, step):
        self.progress_bar.setValue(value)
        self.status_label.setText(status)
        self.step_label.setText(f"단계: {step}")

class ProgressTracker(QObject):
    progress_updated = pyqtSignal(int, str, str)
    status_updated = pyqtSignal(str)

    def __init__(self):
        super().__init__()
        self.current_step = 0
        self.total_steps = 0
        self.is_cancelled = False

    def set_total_steps(self, total):
        self.total_steps = total

    def update_progress(self, step, status):
        self.current_step = step
        progress = int((step / self.total_steps) * 100)
        self.progress_updated.emit(progress, status, f"{step}/{self.total_steps}")

    def update_status(self, status):
        self.status_updated.emit(status)

    def cancel(self):
        self.is_cancelled = True

    def check_cancelled(self):
        if self.is_cancelled:
            raise asyncio.CancelledError("작업이 취소되었습니다.")

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()
        self.current_task = None

    def initUI(self):
        self.setWindowTitle('NF 부동산 물건자료 생성기')
        self.setGeometry(100, 100, 1300, 800)

        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        main_layout = QVBoxLayout(central_widget)

        # 입력 그룹
        input_group = QGroupBox("입력")
        input_layout = QGridLayout()
        input_group.setLayout(input_layout)

        self.lang_combo = QComboBox()
        self.lang_combo.addItems(["한국어", "일본어"])
        input_layout.addWidget(QLabel("언어:"), 0, 0)
        input_layout.addWidget(self.lang_combo, 0, 1)

        self.customer_name_input = QLineEdit()
        input_layout.addWidget(QLabel("고객 이름:"), 1, 0)
        input_layout.addWidget(self.customer_name_input, 1, 1)

        self.property_urls_input = QTextEdit()
        self.property_urls_input.setMinimumHeight(100)
        input_layout.addWidget(QLabel("부동산 URL:"), 2, 0)
        input_layout.addWidget(self.property_urls_input, 2, 1)

        main_layout.addWidget(input_group)

        # 버튼 레이아웃
        button_layout = QHBoxLayout()
        self.start_button = QPushButton("시작")
        self.start_button.clicked.connect(self.start_processing)
        button_layout.addWidget(self.start_button)

        self.cancel_button = QPushButton("취소")
        self.cancel_button.clicked.connect(self.cancel_processing)
        self.cancel_button.setEnabled(False)
        button_layout.addWidget(self.cancel_button)

        main_layout.addLayout(button_layout)

        # 진행 상황 그룹
        progress_group = QGroupBox("진행 상황")
        progress_layout = QVBoxLayout()
        progress_group.setLayout(progress_layout)

        self.progress_widget = ProgressWidget()
        progress_layout.addWidget(self.progress_widget)

        self.detailed_status = QTextEdit()
        self.detailed_status.setReadOnly(True)
        self.detailed_status.setMinimumHeight(150)
        progress_layout.addWidget(self.detailed_status)

        main_layout.addWidget(progress_group)

        # 상태바 생성
        self.statusBar = QStatusBar()
        self.setStatusBar(self.statusBar)

        # 프로그레스 트래커 설정
        self.progress_tracker = ProgressTracker()
        self.progress_tracker.progress_updated.connect(self.update_progress)
        self.progress_tracker.status_updated.connect(self.update_status)

        # 툴팁 설정
        self.setup_tooltips()

        # 도움말 버튼 설정
        self.setup_help_button()

        # 컨텍스트 헬프 모드 설정
        self.setWindowFlags(self.windowFlags() | Qt.WindowType.WindowContextHelpButtonHint)

    def setup_tooltips(self):
        QToolTip.setFont(QFont('SansSerif', 10))
        self.setToolTip('부동산 정보를 처리하여 PPT를 생성합니다.')
        self.lang_combo.setToolTip('출력 언어를 선택하세요.')
        self.customer_name_input.setToolTip('고객의 이름을 입력하세요.')
        self.property_urls_input.setToolTip('처리할 부동산 URL을 한 줄에 하나씩 입력하세요.')
        self.start_button.setToolTip('처리를 시작합니다.')
        self.cancel_button.setToolTip('진행 중인 처리를 취소합니다.')

    def setup_help_button(self):
        help_button = QPushButton("?", self)
        help_button.setFixedSize(20, 20)
        help_button.setToolTip("도움말")
        help_button.clicked.connect(self.show_help)
        self.statusBar.addPermanentWidget(help_button)

    def show_help(self):
        help_text = """
        부동산 정보 처리기 사용 방법:

        1. 언어 선택: 원하는 출력 언어를 선택합니다.
        2. 고객 이름 입력: 고객의 이름을 입력합니다.
        3. URL 입력: 처리할 부동산 정보 페이지의 URL을 입력합니다.
           여러 개의 URL은 줄바꿈으로 구분합니다.
        4. 시작: '시작' 버튼을 클릭하여 처리를 시작합니다.
        5. 진행 상황: 하단의 진행 바와 상태 메시지로 진행 상황을 확인합니다.
        6. 취소: 필요시 '취소' 버튼을 클릭하여 처리를 중단할 수 있습니다.

        특정 기능에 대한 자세한 설명은 해당 요소 위에 마우스를 올려 툴팁을 확인하세요.
        """
        QMessageBox.information(self, "도움말", help_text)

    def eventFilter(self, obj, event):
        if event.type() == QEvent.Type.WhatsThis:
            QWhatsThis.enterWhatsThisMode()
            return True
        return super().eventFilter(obj, event)

    @asyncSlot()
    async def start_processing(self):
        lang = self.lang_combo.currentText()
        customer_name = self.customer_name_input.text()
        property_urls = self.property_urls_input.toPlainText().strip().split("\n")

        if not customer_name or not property_urls:
            QMessageBox.critical(self, "Error", "모든 필드를 입력해주세요.")
            return

        lang_code = 'ko' if lang == '한국어' else 'ja'
        self.detailed_status.clear()
        self.progress_widget.update_progress(0, "처리 시작...", "0/0")
        self.start_button.setEnabled(False)
        self.cancel_button.setEnabled(True)
        
        self.progress_tracker.is_cancelled = False
        
        try:
            self.current_task = asyncio.create_task(main(property_urls, customer_name, lang_code, self.progress_tracker))
            await self.current_task
        except asyncio.CancelledError:
            self.update_status("작업이 취소되었습니다.")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"처리 중 오류가 발생했습니다: {str(e)}")
        finally:
            self.start_button.setEnabled(True)
            self.cancel_button.setEnabled(False)
            self.current_task = None

    def cancel_processing(self):
        if self.current_task:
            self.progress_tracker.cancel()

    def update_progress(self, value, status, step):
        self.progress_widget.update_progress(value, status, step)

    def update_status(self, status):
        self.detailed_status.append(status)
        self.detailed_status.verticalScrollBar().setValue(self.detailed_status.verticalScrollBar().maximum())

async def process_property(prs, property_url, slide_index, customer_name, translator, lang_code, progress_tracker, temp_dir):
    try:
        progress_tracker.check_cancelled()
        progress_tracker.update_progress(slide_index * 5 - 4, "부동산 정보 스크래핑 중...")
        property_info, total_floors = await scrape_property_info(property_url)

        progress_tracker.check_cancelled()
        progress_tracker.update_progress(slide_index * 5 - 3, "지도 이미지 다운로드 중...")
        address = property_info.address
        coordinates = get_coordinates(address)
        map_image = None
        if coordinates:
            x, y = coordinates
            map_url = get_static_map_url(x, y)
            map_image = os.path.join(temp_dir, f"map_{property_info.complex_name}.png")
            download_image(map_url, map_image)

        progress_tracker.check_cancelled()
        progress_tracker.update_progress(slide_index * 5 - 2, "평면도 다운로드 중...")
        floor_plan_image = await download_floor_plan(property_url, temp_dir)

        progress_tracker.check_cancelled()
        progress_tracker.update_progress(slide_index * 5 - 1, "단지 사진 다운로드 중...")
        complex_photo = await download_complex_photo(property_url, temp_dir)

        if slide_index > 1:
            new_slide = add_new_slide(prs)
            slide_index = len(prs.slides) - 1

        if lang_code == 'ja':
            await property_info.translate(translator)

        progress_tracker.check_cancelled()
        progress_tracker.update_progress(slide_index * 5, "PowerPoint 슬라이드 업데이트 중...")
        update_pptx_template(prs, property_info, total_floors, customer_name, slide_index, lang_code)

        image_map = {
            14: complex_photo,
            21: floor_plan_image,
            26: map_image,
            9: resource_path("회사로고.png"),
            10: resource_path("글로벌부동산.jpg")
        }

        add_images_to_rectangles(prs, image_map, slide_index)
        progress_tracker.update_status(f"완료: {property_url}")

        # 중개사 정보 추가
        broker_info = await scrape_broker_info(property_url)
        add_slide_note_with_broker_info(prs, slide_index, broker_info)

    except asyncio.CancelledError:
        raise
    except Exception as e:
        progress_tracker.update_status(f"오류 발생: {property_url} - {str(e)}")

async def main(property_urls, customer_name, lang_code, progress_tracker):
    # 임시 디렉토리 생성
    with tempfile.TemporaryDirectory() as temp_dir:
        template_path = resource_path("template.pptx") if lang_code == 'ko' else resource_path("template_ja.pptx")
        prs = Presentation(template_path)
        update_cover_slide(prs, customer_name)
        translator = DeepLTranslator() if lang_code == 'ja' else None

        progress_tracker.set_total_steps(len(property_urls) * 5)

        for index, url in enumerate(property_urls, start=1):
            try:
                progress_tracker.check_cancelled()
                await process_property(prs, url, index, customer_name, translator, lang_code, progress_tracker, temp_dir)
            except asyncio.CancelledError:
                raise
            except Exception as e:
                progress_tracker.update_status(f"오류 발생: 매물 {index} - {str(e)}")
                continue

        output_path = "updated_template.pptx"
        prs.save(output_path)
        progress_tracker.update_status("처리 완료. PowerPoint 파일이 저장되었습니다.")

    # 임시 디렉토리는 with 문을 벗어나면 자동으로 삭제됩니다.

def install_playwright():
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            # 모든 브라우저 설치
            p.chromium.install()
            p.firefox.install()
            p.webkit.install()
    except Exception as e:
        print(f"Playwright 설치 중 오류 발생: {e}")

if __name__ == '__main__':
    try:
        # Playwright 브라우저 설치 경로 설정
        import os
        import sys
        
        # 실행 파일 디렉토리 확인
        if getattr(sys, 'frozen', False):
            application_path = os.path.dirname(sys.executable)
        else:
            application_path = os.path.dirname(os.path.abspath(__file__))
            
        # Playwright 브라우저 경로 설정
        os.environ["PLAYWRIGHT_BROWSERS_PATH"] = os.path.join(application_path, "playwright-browsers")
        
        # Playwright 브라우저 설치
        install_playwright()
        
        app = QApplication(sys.argv)
        app.setStyle('WindowsVista')
        loop = QEventLoop(app)
        asyncio.set_event_loop(loop)
        
        main_window = MainWindow()
        main_window.show()
        
        with loop:
            loop.run_forever()
    except Exception as e:
        print(f"An error occurred: {e}")
        input("Press Enter to exit...")
        raise