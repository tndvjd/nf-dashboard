from pptx import Presentation

def analyze_pptx_template(pptx_path):
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        print(f"Slide Layout: {slide.slide_layout.name}")
        
        for shape_id, shape in enumerate(slide.shapes):
            print(f"\nShape ID: {shape_id}")
            print(f"Shape Type: {type(shape).__name__}")
            print(f"Shape Name: {shape.name}")
            print(f"Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}")
            
            if hasattr(shape, 'table'):
                table = shape.table
                print(f"Table with {len(table.rows)} rows and {len(table.columns)} columns")
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        print(f"  Row {row_idx}, Col {col_idx}: {cell.text}")
            
            elif hasattr(shape, 'image'):
                print("This shape contains an image")
            
            elif hasattr(shape, 'text'):
                print(f"Text: {shape.text}")
            
            # 이미지 컨테이너로 사용될 수 있는 도형 식별
            if hasattr(shape, 'fill'):
                print("This shape might be used as an image container")

if __name__ == "__main__":
    template_path = "template.pptx"  # 여기에 실제 템플릿 파일 경로를 입력하세요
    analyze_pptx_template(template_path)