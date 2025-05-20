# ppt_utils.py

import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict
from scraper import PropertyInfo
import json


def update_cell(cell, text: str):
    """
    PowerPoint 테이블 셀의 텍스트를 업데이트하고 스타일을 설정합니다.

    Args:
        cell: PowerPoint 테이블 셀 객체
        text (str): 셀에 입력할 텍스트
    """
    cell.text = text
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 텍스트 중앙 정렬
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)  # 폰트 크기 설정
            run.font.color.rgb = RGBColor(0, 0, 0)  # 폰트 색상 설정 (검정색)
            run.font.bold = False  # 볼드 해제


def add_new_slide(prs: Presentation):
    """
    PowerPoint 프레젠테이션에 새로운 슬라이드를 추가합니다.

    Args:
        prs (Presentation): PowerPoint 프레젠테이션 객체

    Returns:
        pptx.slide.Slide: 새로 추가된 슬라이드 객체
    """
    source_slide = prs.slides[1]  # 두 번째 슬라이드를 템플릿으로 사용
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # 모든 도형을 새 슬라이드로 복사
    for shape in source_slide.shapes:
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Title과 Content Placeholder 제거
    shapes_to_remove = []
    for shape in new_slide.shapes:
        if shape.name.startswith('Title') or shape.name.startswith('Content Placeholder'):
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    print(f"Added new slide. Total slides: {len(prs.slides)}")
    return new_slide


def add_images_to_rectangles(prs: Presentation, image_map: Dict[int, str], slide_index: int):
    """
    PowerPoint 슬라이드의 특정 도형에 이미지를 추가합니다.

    Args:
        prs (Presentation): PowerPoint 프레젠테이션 객체
        image_map (Dict[int, str]): 도형 ID와 이미지 경로를 매핑하는 딕셔너리
        slide_index (int): 이미지를 추가할 슬라이드 인덱스
    """
    print(f"Adding images to slide {slide_index}")
    slide = prs.slides[slide_index]

    # 기존 이미지 제거 (회사로고와 글로벌부동산 제외)
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.shape_id not in [9, 10]:  # 회사로고, 글로벌부동산 제외
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in slide.shapes:
        if shape.shape_id in image_map:
            image_path = image_map[shape.shape_id]
            if image_path:
                if shape.shape_id in [9, 10]:  # 회사로고와 글로벌부동산
                    # 원래 크기 그대로 유지
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                else:
                    # 도형의 크기에서 5px(2.5px * 2)를 뺀 크기로 이미지 삽입
                    left = shape.left + Pt(2.5)
                    top = shape.top + Pt(2.5)
                    width = shape.width - Pt(5)
                    height = shape.height - Pt(5)
                
                slide.shapes.add_picture(image_path, left, top, width, height)
                print(f"이미지 추가 완료: {image_path}, 도형 ID: {shape.shape_id}")


def sqm_to_pyeong(sqm: float) -> float:
    """
    제곱미터를 평으로 변환합니다.

    Args:
        sqm (float): 제곱미터 값

    Returns:
        float: 평 값
    """
    return round(sqm / 3.3058, 2)


def format_area_with_pyeong(area: str, lang: str = 'ko') -> str:
    """
    면적 정보를 '제곱미터/평' 형식으로 변환합니다.

    Args:
        area (str): 면적 정보 문자열
        lang (str): 언어 코드 ('ko' 또는 'ja')

    Returns:
        str: '제곱미터/평' 형식으로 변환된 문자열
    """
    match = re.search(r'([\d.]+)㎡', area)
    if match:
        sqm = float(match.group(1))
        pyeong = sqm_to_pyeong(sqm)
        if lang == 'ja':
            return f"{sqm}㎡/{pyeong}坪"  # 일본어로 '평'을 '坪'으로 변환
        return f"{sqm}㎡/{pyeong}평"  # 기본값은 한국어
    return area


def clean_complex_name(name: str) -> str:
    """
    단지 이름에서 괄호 안의 내용을 제거합니다.

    Args:
        name (str): 단지 이름 문자열

    Returns:
        str: 괄호 안의 내용이 제거된 단지 이름 문자열
    """
    return re.sub(r'\([^)]*\)$', '', name).strip()


def update_property_info(table, property_info: PropertyInfo, lang: str = 'ko'):
    """
    PowerPoint 슬라이드의 '물건정보표' 테이블을 업데이트합니다.

    Args:
        table: PowerPoint 테이블 객체
        property_info (PropertyInfo): 매물 정보 객체
        lang (str): 언어 코드 ('ko' 또는 'ja')
    """
    print("Updating property info:")
    update_cell(table.cell(0, 1), property_info.floor)
    print(f"Floor: {property_info.floor}")

    areas = property_info.area.split('/')
    if len(areas) == 2:
        contract_area = format_area_with_pyeong(areas[0], lang)
        exclusive_area = format_area_with_pyeong(areas[1], lang)
        update_cell(table.cell(1, 1), contract_area)
        update_cell(table.cell(2, 1), exclusive_area)
        print(f"Contract area: {contract_area}")
        print(f"Exclusive area: {exclusive_area}")

    update_cell(table.cell(3, 1), property_info.rooms)
    print(f"Rooms: {property_info.rooms}")
    update_cell(table.cell(4, 1), property_info.price)
    print(f"Price: {property_info.price}")
    update_cell(table.cell(5, 1), property_info.maintenance_fee)
    print(f"Maintenance fee: {property_info.maintenance_fee}")
    update_cell(table.cell(6, 1), property_info.direction)
    print(f"Direction: {property_info.direction}")
    update_cell(table.cell(7, 1), property_info.move_in_date)
    print(f"Move-in date: {property_info.move_in_date}")


def update_complex_info(table, property_info: PropertyInfo, total_floors: str, lang: str = 'ko'):
    """
    PowerPoint 슬라이드의 '단지정보표' 테이블을 업데이트합니다.
    한국어와 일본어를 모두 지원합니다.

    Args:
        table: PowerPoint 테이블 객체
        property_info (PropertyInfo): 매물 정보 객체
        total_floors (str): 총 층수 정보
        lang (str): 언어 코드 ('ko' 또는 'ja')
    """
    print("Updating complex info:")
    for row in table.rows:
        cell_text = row.cells[0].text
        if cell_text in ["소재지", "所在地"]:
            update_cell(row.cells[1], property_info.address)
            print(f"Address/所在地: {property_info.address}")
        elif cell_text in ["건축연도", "建築年度"]:
            update_cell(row.cells[1], property_info.approval_date)
            print(f"Approval date/建築年度: {property_info.approval_date}")
        elif cell_text in ["세대수", "世帯数"]:
            update_cell(row.cells[1], property_info.total_units)
            print(f"Total units/世帯数: {property_info.total_units}")
        elif cell_text in ["총층수", "総階数"]:
            if lang == 'ja':
                update_cell(row.cells[1], f"{property_info.total_floors}階")  # 일본어로 '층'을 '階'로 변환
            else:
                update_cell(row.cells[1], f"{property_info.total_floors}층")
            print(f"Total floors/総階数: {property_info.total_floors}층")
        elif cell_text in ["난방", "暖房"]:
            update_cell(row.cells[1], property_info.heating)
            print(f"Heating/暖房: {property_info.heating}")


def update_cover_slide(prs: Presentation, customer_name: str):
    """
    PowerPoint 템플릿의 표지 슬라이드를 업데이트합니다.
    """
    print("Updating cover slide")
    slide = prs.slides[0]  # 첫 번째 슬라이드 (표지)
    for shape in slide.shapes:
        print(f"Shape ID: {shape.shape_id}, Shape Name: {shape.name}")  # Shape ID와 이름 출력
        if shape.shape_id == 8:  # Shape ID 2 확인
            print("Target shape found!")
            if shape.has_table:
                print("Shape has table.")
                table = shape.table
                print(f"Table rows: {len(table.rows)}, Table columns: {len(table.columns)}")
                if len(table.rows) > 0 and len(table.columns) > 1:
                    cell = table.cell(0, 1)
                    print(f"Original cell text: {cell.text}")  # 원래 셀 텍스트 출력
                    cell.text = customer_name
                    print(f"Updated customer name: {customer_name}")
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(14)
                            run.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
                            run.font.bold = True  # 글씨 진하게 처리
            else:
                print("Shape does not have table.")  # 테이블이 없는 경우 출력


def update_property_slide(prs: Presentation, property_info: PropertyInfo, total_floors: str, customer_name: str, slide_index: int, lang: str = 'ko'):
    """
    PowerPoint 템플릿의 매물 정보 슬라이드를 업데이트합니다.
    """
    try:
        print(f"Updating property info slide {slide_index}")
        if slide_index >= len(prs.slides):
            print(f"Error: Slide index {slide_index} is out of range. Total slides: {len(prs.slides)}")
            return

        slide = prs.slides[slide_index]
        district = f"[{property_info.district}]" if property_info.district else ""

        for shape in slide.shapes:
            if shape.name == "물건명":
                cleaned_name = clean_complex_name(property_info.complex_name)
                text = f"No.{slide_index} {district} {cleaned_name} ({property_info.complex_type})"
                shape.text_frame.text = text
                print(f"Updated property name: {text}")
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT  # 우측 정렬
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
                        run.font.bold = False  # Bold 해제
            elif shape.has_table:
                if shape.name == "물건정보표":
                    print("Updating property info table")
                    update_property_info(shape.table, property_info, lang)
                elif shape.name == "단지정보표":
                    print("Updating complex info table")
                    update_complex_info(shape.table, property_info, total_floors, lang)
    except Exception as e:
        print(f"Error in update_property_slide: {e}")


def save_to_json(property_info: PropertyInfo, total_floors: str, customer_name: str, slide_index: int):
    """
    정제된 부동산 정보를 JSON 파일로 저장합니다.

    Args:
        property_info (PropertyInfo): 부동산 정보 객체
        total_floors (str): 총 층수 정보
        customer_name (str): 고객 이름
        slide_index (int): 슬라이드 인덱스
    """
    # 면적 정보를 분리하여 저장
    areas = property_info.area.split('/')
    contract_area = format_area_with_pyeong(areas[0]) if len(areas) > 0 else ""
    exclusive_area = format_area_with_pyeong(areas[1]) if len(areas) > 1 else ""

    data = {
        "slide_index": slide_index,
        "customer_name": customer_name,
        "property_info": {
            "contract_area": contract_area,
            "exclusive_area": exclusive_area,
            "floor": property_info.floor,
            "rooms": property_info.rooms,
            "maintenance_fee": property_info.maintenance_fee,
            "direction": property_info.direction,
            "move_in_date": property_info.move_in_date,
            "address": property_info.address,
            "total_units": property_info.total_units,
            "approval_date": property_info.approval_date,
            "heating": property_info.heating,
            "complex_name": property_info.complex_name,
            "complex_type": property_info.complex_type,
            "price": property_info.price
        },
        "total_floors": f"{total_floors}층"  # "층" 추가
    }

    json_filename = f"property_info_{property_info.complex_name}_{slide_index}.json"
    with open(json_filename, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
    print(f"Property info saved to {json_filename}")

def set_global_font(prs, lang='ko', font_size=Pt(12)):
    """
    프레젠테이션의 모든 텍스트에 대해 글꼴과 크기를 설정합니다.

    Args:
        prs (Presentation): PowerPoint 프레젠테이션 객체
        lang (str): 언어 코드 ('ko' 또는 'ja')
        font_size (Pt): 설정할 글꼴 크기
    """
    font_name = '맑은 고딕' if lang == 'ko' else 'Meiryo UI'
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = font_size
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = font_name
                                run.font.size = font_size


def update_pptx_template(prs: Presentation, property_info: PropertyInfo, total_floors: str, customer_name: str, slide_index: int, lang: str = 'ko'):
    """
    PowerPoint 템플릿을 업데이트합니다.
    """
    print(f"Updating template for slide {slide_index}")
    print(f"Current number of slides: {len(prs.slides)}")

    if slide_index == 0:  # 첫 번째 슬라이드 (표지)
        update_cover_slide(prs, customer_name)
    else:  # 부동산 정보 슬라이드
        update_property_slide(prs, property_info, total_floors, customer_name, slide_index, lang)

    # 프레젠테이션의 모든 텍스트에 대해 글꼴 설정
    set_global_font(prs, lang=lang, font_size=Pt(12))

    print(f"Finished updating template for slide {slide_index}")

def add_slide_note_with_broker_info(prs: Presentation, slide_index: int, broker_info: str):
    """
    슬라이드에 중개사 정보를 메모로 추가합니다.

    Args:
        prs (Presentation): PowerPoint 프레젠테이션 객체
        slide_index (int): 슬라이드 인덱스
        broker_info (str): 추가할 중개사 정보
    """
    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = broker_info