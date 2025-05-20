from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE # MSO_SHAPE_TYPE 자체를 사용
from pptx.util import Inches, Pt
import datetime
import math
# import requests # 실제 이미지 로드 시 주석 해제
# from io import BytesIO # 실제 이미지 로드 시 주석 해제

TEMPLATE_PATH = "ppt_template.pptx"
OUTPUT_PATH_DEBUG = "최종_생성된_물건자료.pptx"
IMAGE_BASE_URL = "https://land.naver.com" # 실제 이미지 경로의 기본 URL

# --- 1. 샘플 데이터 정의 (이전과 동일 또는 약간 수정) ---
client_and_document_data = {
    "문서명": "정식 사택 물건 자료", # 매핑 데이터의 '텍스트_문서명'에 해당
    "고객명": "신지성 고객님",     # 매핑 데이터의 '텍스트_고객명'에 해당
    "회사명": "AGC일렉트로닉스코리아 (주)" # 매핑 데이터의 '텍스트_회사명'에 해당
}

# 매물 데이터 정의 - 여러 개의 매물을 테스트하기 위한 데이터
# 첫 번째 매물데이터
first_property_data = { # 이전 JSON과 동일한 구조로 가정, 필요한 값 위주로 다시 정리
    "articleDetail": {
        "articleNo": "2523213850", "articleName": "보타니끼논현 1동", "divisionName": "강남구",
        "aptName": "보타니끼논현", "exposureAddress": "서울시 강남구 논현동 123-45번지", # 상세 주소 추가
        "aptUseApproveYmd": "20241122", "aptHouseholdCount": "42",
        "aptHeatMethodTypeName": "개별난방", "aptHeatFuelTypeName": "도시가스",
        "buildingName": "101동", "roomCount": "1", "bathroomCount": "1", # '동' 정보 수정
        "moveInTypeName": "즉시입주",
        "articleFeatureDescription": "강남 중심, 즉시입주 가능한 최고급 하이엔드 오피스텔입니다.",
        "detailDescription": "사진은 실제 방문하여 촬영한 것이며, 보증금 및 월세 조건은 조율 가능합니다. 24시간 상담 환영!\n복층형 구조에 넓은 테라스 보유. 분리형 원룸 및 화장실 1개. 주차는 세대당 1.45대 가능(총 61대).",
        "tagList": ["2년이내신축", "테라스굷", "역세권도보5분", "화장실1개"],
        "grandPlanList": [{"imageSrc": "/20220706_249/hscp_img_1657095670067Nphxa_JPEG/photoinfra_1657095669504.jpg", "imageType": "14"}], # 평면도/단지배치도
        "latitude": "37.51189", "longitude": "127.035282",
    },
    "articleAddition": {
        "floorInfo": "중층/9", "direction": "남동향", # '중층' 정보 추가, 방향 수정
        "dealOrWarrantPrc": "1억", "rentPrc": "500",
        "representativeImgUrl": "/20250508_24/1746696666335n06Ea_JPEG/a80fd2f71dc5fe0c716ffda69d867e49.JPG"
    },
    "articleFloor": {"totalFloorCount": "9", "buildingHighestFloor": "9"}, # 실제로는 totalFloorCount가 단지 전체 최고층, buildingHighestFloor가 해당 동 최고층일 수 있음
    "articlePrice": {"warrantPrice": 10000, "rentPrice": 500}, # 만원 단위
    "articleSpace": {"supplySpace": 111.88, "exclusiveSpace": 47.35}, # ㎡ 단위
    "articlePhotos": [ # 대표 이미지 외 추가 사진들
        {"imageSrc": "/20250508_198/1746696668451MhMlb_JPEG/KakaoTalk_20250430_143429280_09.jpg", "imageType": "10"},
        {"imageSrc": "/20250508_212/1746696668451y2OrJ_JPEG/KakaoTalk_20250430_143429280_01.jpg", "imageType": "10"},
    ],
    "administrationCostInfo": {"chargeCodeType": "03"}, # 관리비 금액 정보 없음
    # "별도 입력 필요" 항목들 (FastAPI로 받을 예정이지만 테스트 위해 임시값)
    "참고사항_입력": "6월 중 입주 협의 가능합니다. 반려동물 동반은 어려운니다.",
    "비고_입력": "현관 보안 철저, 건물 내 편의시설(피트니스, 라운지) 이용 가능."
}

# 두 번째 매물데이터 (테스트용)
second_property_data = {
    "articleDetail": {
        "articleNo": "2523123456", "articleName": "래미안 퍼스티지", "divisionName": "강남구",
        "aptName": "래미안 퍼스티지", "exposureAddress": "서울시 강남구 신사동 123-45", 
        "aptUseApproveYmd": "20150312", "aptHouseholdCount": "320",
        "aptHeatMethodTypeName": "개별난방", "aptHeatFuelTypeName": "도시가스",
        "buildingName": "101동", "roomCount": "3", "bathroomCount": "2",
        "moveInTypeName": "즉시 입주 가능",
        "articleFeatureDescription": "신사역 도보 5분 거리, 리모델링 완료된 남향 코너 호실",
        "detailDescription": "채광 좋은 거실, 입주 즉시 가능한 깔끔한 아파트입니다.",
        "tagList": ["남향", "역세권", "반려동물 가능", "신규 등록"],
        "grandPlanList": [{"imageSrc": "/타입1_평면도.jpg", "imageType": "14"}],
        "latitude": "37.5167", "longitude": "127.0402"
    },
    "articleAddition": {
        "floorInfo": "12/15", "direction": "남향",
        "dealOrWarrantPrc": "25000", "rentPrc": "120",
        "representativeImgUrl": "/매물사진2.jpg"
    },
    "articleFloor": {
        "totalFloorCount": "15", "buildingHighestFloor": "15"
    },
    "articlePrice": {"warrantPrice": 25000, "rentPrice": 120},
    "articleSpace": {"supplySpace": 112.32, "exclusiveSpace": 84.83},
    "articlePhotos": [
        {"imageSrc": "/매물사진2_1.jpg", "imageType": "INTERNAL"},
        {"imageSrc": "/매물사진2_2.jpg", "imageType": "INTERNAL"}
    ],
    "administrationCostInfo": {"chargeCodeType": "INCLUDE_ALL"},
    "참고사항_입력": "관리비에 수도/전기/가스 포함",
    "비고_입력": "주차 2대 가능"
}

# 세 번째 매물데이터 (테스트용)
third_property_data = {
    "articleDetail": {
        "articleNo": "2523654321", "articleName": "센트럴파크 타워", "divisionName": "강남구",
        "aptName": "센트럴파크 타워", "exposureAddress": "서울시 강남구 역삼동 789-10", 
        "aptUseApproveYmd": "20180525", "aptHouseholdCount": "192",
        "aptHeatMethodTypeName": "중앙난방", "aptHeatFuelTypeName": "지역난방",
        "buildingName": "A동", "roomCount": "2", "bathroomCount": "1",
        "moveInTypeName": "2개월 후 입주 가능",
        "articleFeatureDescription": "역삼역 초역세권, 풀옵션 신축급 오피스텔",
        "detailDescription": "2년 계약 시 중개수수료 할인 가능합니다.",
        "tagList": ["풀옵션", "신축", "주차가능", "24시간 보안"],
        "grandPlanList": [{"imageSrc": "/타입2_평면도.jpg", "imageType": "14"}],
        "latitude": "37.5006", "longitude": "127.0359"
    },
    "articleAddition": {
        "floorInfo": "8/18", "direction": "동향",
        "dealOrWarrantPrc": "15000", "rentPrc": "90",
        "representativeImgUrl": "/매물사진3.jpg"
    },
    "articleFloor": {
        "totalFloorCount": "18", "buildingHighestFloor": "18"
    },
    "articlePrice": {"warrantPrice": 15000, "rentPrice": 90},
    "articleSpace": {"supplySpace": 63.6, "exclusiveSpace": 45.5},
    "articlePhotos": [
        {"imageSrc": "/매물사진3_1.jpg", "imageType": "INTERNAL"},
        {"imageSrc": "/매물사진3_2.jpg", "imageType": "INTERNAL"}
    ],
    "administrationCostInfo": {"chargeCodeType": "INCLUDE_PART"},
    "참고사항_입력": "보안키 및 디지털 도어락 설치",
    "비고_입력": "공용 사우나, 피트니스센터 이용 가능"
}

# --- 2. 매핑 규칙 정의 (디버깅 결과 및 사용자 CSV 기반) ---
# 키: PPTX 템플릿의 "선택 창"에 나오는 실제 도형 이름
# 값: 튜플 (자리표시자텍스트_또는_테이블셀좌표, JSON경로/키_또는_리스트, 가공함수이름, 후처리텍스트_또는_폰트정보)
# 가공함수 이름이 None이면 직접 값을 사용
# 테이블 셀 좌표는 (행, 열) 튜플. 값에 "TABLE_HEADER" 등이면 해당 셀은 헤더로 간주.
mappings = {
    # 슬라이드 1
    "txt_document_title": ("{{document_title}}", "문서명", None, {"font_name": "나눔고딕", "font_size": Pt(24)}),
    "txt_client_name": ((0,0), "고객명", None, {"font_name": "나눔고딕", "font_size": Pt(18)}), # 테이블 셀
    "txt_company_name": ("{{company_name}}", "회사명", None, {"font_name": "나눔고딕", "font_size": Pt(16)}),

    # 슬라이드 2
    # 매물 타이틀 (txt_property_title TABLE (0,0) 셀)
    "txt_property_title": ((0,0), ["articleDetail.divisionName", "articleDetail.aptName"], 'format_property_title_text', {"font_name": "나눔고딕", "font_size": Pt(18)}),

    # 매물 순번 표시 (TEXT_BOX)
    "txt_property_page": ("{{매물순번_표시}}", None, 'get_property_order', {"font_name": "나눔고딕", "font_size": Pt(10)}),

    # 단지 정보 테이블 (tbl_complex_info) - 5행 2열
    # (도형이름, (행,열), JSON경로, 가공함수, {폰트정보})
    "tbl_complex_info": {
        (0, 0): ("단지주소", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, "articleDetail.exposureAddress", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("준공연도", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleDetail.aptUseApproveYmd", 'format_date_yyyy_mm', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("총세대수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.aptHouseholdCount", 'format_household_count', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("총층수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, "articleFloor.buildingHighestFloor", 'format_floor_count', {"font_name": "나눔고딕", "font_size": Pt(10)}), # buildingHighestFloor 또는 totalFloorCount
        (4, 0): ("난방방식", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.aptHeatMethodTypeName", "articleDetail.aptHeatFuelTypeName"], 'format_heating_info', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },

    # 물건 정보 테이블 1 (tbl_property_detail_1) - 5행 2열
    "tbl_property_detail_1": {
        (0, 0): ("동·호수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articleDetail.buildingName", "articleAddition.floorInfo"], 'format_dong_ho_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("계약면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleSpace.supplySpace", 'format_supply_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("전용면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleSpace.exclusiveSpace", 'format_exclusive_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("방수/욕실수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.roomCount", "articleDetail.bathroomCount"], 'format_room_bathroom_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("방향", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, "articleAddition.direction", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },

    # 물건 정보 테이블 2 (tbl_property_detail_2) - 5행 2열
    "tbl_property_detail_2": {
        (0, 0): ("보증금/월세", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articlePrice.warrantPrice", "articlePrice.rentPrice", "articleAddition.dealOrWarrantPrc", "articleAddition.rentPrc"], 'format_price_details_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("기본관리비", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "administrationCostInfo.chargeCodeType", 'format_management_fee_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("입주가능일", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.moveInTypeName", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("참고사항", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.articleFeatureDescription", "참고사항_입력"], 'get_first_available_value', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("비고", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.tagList", "articleDetail.detailDescription", "비고_입력"], 'format_note_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },

    # 이미지 도형 이름
    "img_complex_view": (None, "articleAddition.representativeImgUrl", 'prepare_image_url', None),
    "img_complex_floorplan": (None, "articleDetail.grandPlanList[0].imageSrc", 'prepare_image_url_if_type14', None), # grandPlanList가 비어있을 수 있으니 주의
    "img_complex_mapimg": (None, ["articleDetail.latitude", "articleDetail.longitude"], 'generate_map_image_placeholder', None),
}

# --- 3. 데이터 가공 헬퍼 함수들 ---
def get_nested_value(data_dict, keys_string_or_list, default=""):
    if keys_string_or_list is None: return default
    if isinstance(keys_string_or_list, list):
        for keys_string in keys_string_or_list:
            val = get_single_nested_value(data_dict, keys_string, None)
            if val is not None and val != "": return val
        return default
    return get_single_nested_value(data_dict, keys_string_or_list, default)

def get_single_nested_value(data_dict, keys_string, default=""):
    if keys_string is None:
        return default
    keys = keys_string.split('.')
    value = data_dict
    try:
        for key_part in keys:
            if '[' in key_part and key_part.endswith(']'):
                key = key_part.split('[')[0]
                index = int(key_part.split('[')[1][:-1])
                if key in value and isinstance(value[key], list) and len(value[key]) > index:
                    value = value[key][index]
                else: return default
            elif isinstance(value, dict) and key_part in value:
                value = value[key_part]
            else: return default
        return value if value is not None else default
    except (KeyError, IndexError, TypeError, AttributeError): return default


def format_date_yyyy_mm(value, data_dict=None, json_path=None):
    if not value or len(str(value)) != 8: return ""
    try:
        dt_obj = datetime.datetime.strptime(str(value), "%Y%m%d")
        return f"{dt_obj.year}년 {dt_obj.month}월"
    except ValueError: return ""

def calculate_pyeong(m2_value, precision=1): # 소수점 한자리로 변경
    if isinstance(m2_value, (int, float)) and m2_value > 0:
        return round(m2_value / 3.3058, precision)
    return ""

def format_property_title_text(values, data_dict=None, json_path=None):
    # values는 [divisionName, aptName] 리스트
    prop_order = "1" # 여기서 매물 순번을 어떻게 가져올지 결정 필요 (예: 외부 변수 또는 추가 인자)
    div_name = values[0] if len(values) > 0 and values[0] else ""
    apt_name = values[1] if len(values) > 1 and values[1] else ""
    return f"No.{prop_order} [{div_name}] {apt_name}"

def get_property_order(value=None, data_dict=None, json_path=None):
    return "1" # 단순 반환

def format_household_count(value, data_dict=None, json_path=None):
    return f"{value}세대" if value else ""

def format_floor_count(value, data_dict=None, json_path=None):
    return f"{value}층" if value else ""

def format_heating_info(values, data_dict=None, json_path=None):
    # values는 [aptHeatMethodTypeName, aptHeatFuelTypeName]
    method = values[0] if len(values) > 0 and values[0] else ""
    fuel = values[1] if len(values) > 1 and values[1] else ""
    if method and fuel: return f"{method}, {fuel}"
    return method or fuel or ""

def format_dong_ho_for_ppt(values, data_dict=None, json_path=None):
    # values는 [buildingName, floorInfo]
    building_name = values[0] if len(values) > 0 and values[0] else ""
    floor_info_raw = values[1] if len(values) > 1 and values[1] else ""
    
    floor_display = ""
    if isinstance(floor_info_raw, str):
        if "중층" in floor_info_raw or "고층" in floor_info_raw or "저층" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0]
        elif "/" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0] + "층"
        else: floor_display = floor_info_raw
    else: floor_display = str(floor_info_raw)
    
    return f"{building_name} {floor_display}".strip()

def format_supply_area_text(supply_m2_val, data_dict=None, json_path=None):
    supply_py_val = calculate_pyeong(supply_m2_val)
    return f"{supply_m2_val}㎡ ({supply_py_val}평)" if supply_m2_val else ""

def format_exclusive_area_text(exclusive_m2_val, data_dict=None, json_path=None):
    exclusive_py_val = calculate_pyeong(exclusive_m2_val)
    return f"{exclusive_m2_val}㎡ ({exclusive_py_val}평)" if exclusive_m2_val else ""

def format_room_bathroom_text(values, data_dict=None, json_path=None):
    # values는 [roomCount, bathroomCount]
    room = values[0] if len(values) > 0 and values[0] else "0"
    bath = values[1] if len(values) > 1 and values[1] else "0"
    return f"{room} / {bath}"

def format_management_fee_for_ppt(value, data_dict=None, json_path=None):
    # value는 administrationCostInfo.chargeCodeType
    # 실제 금액 정보가 JSON에 없으므로 하드코딩 또는 다른 로직 필요
    return "확인 어려움" # 또는 "별도 문의"

def format_price_details_for_ppt(values, data_dict=None, json_path=None):
    # values는 [warrantPrice, rentPrice, dealOrWarrantPrc, rentPrc]
    w_price_num = values[0] if len(values) > 0 and values[0] else 0
    r_price_num = values[1] if len(values) > 1 and values[1] else 0
    # deal_text = values[2] if len(values) > 2 and values[2] else "" # dealOrWarrantPrc
    # rent_text = values[3] if len(values) > 3 and values[3] else "" # rentPrc

    # 숫자 기반으로 "X억 Y만원 / Z만원" 형태 구성
    formatted_w = ""
    if isinstance(w_price_num, (int, float)) and w_price_num > 0:
        if w_price_num >= 10000:
            eok = int(w_price_num // 10000)
            manwon_part = int(w_price_num % 10000)
            formatted_w = f"{eok}억"
            if manwon_part > 0: formatted_w += f" {manwon_part:,}"
        else: formatted_w = f"{int(w_price_num):,}"
    
    formatted_r = ""
    if isinstance(r_price_num, (int, float)) and r_price_num > 0:
        formatted_r = f"{int(r_price_num):,}"
        
    if formatted_w and formatted_r: return f"{formatted_w} / {formatted_r} (만원)"
    if formatted_w: return f"{formatted_w} (만원)"
    if formatted_r: return f" / {formatted_r} (만원)" # 전세 없이 월세만
    return "가격 정보 없음"


def get_first_available_value(values, data_dict=None, json_path=None):
    # values는 [json_field1, json_field2_from_client_data, ...]
    # data_dict는 article_json_data임. client_data는 여기서 직접 접근 어려우므로,
    # client_data의 값은 이미 get_nested_value 단계에서 추출되어 values 리스트에 포함되어야 함.
    for val_candidate_key in values: # values의 각 요소는 실제 값이 아니라 JSON 키 문자열이어야 함.
        # 또는, values에 이미 추출된 값이 들어온다고 가정하고 로직 변경
        # 현재 get_nested_value는 values가 키 리스트일 때 첫번째 유효값을 반환함.
        # 이 함수는 그 결과(단일 값)를 받거나, 또는 여러 값의 리스트를 받도록 설계해야 함.
        # 여기서는 values가 이미 추출된 값들의 리스트라고 가정.
        val = val_candidate_key # 만약 values가 ["값1", "값2"] 형태라면
        if val and val != "정보 없음": return str(val)
    return "내용 없음"


def format_note_for_ppt(values, data_dict=None, json_path=None):
    # values는 [tagList, detailDescription, 비고_입력(client_data)]
    tags = values[0] if len(values) > 0 and isinstance(values[0], list) else []
    detail_desc = values[1] if len(values) > 1 and values[1] else ""
    manual_note = values[2] if len(values) > 2 and values[2] else ""

    if manual_note and manual_note != "정보 없음": return manual_note
    if tags: return ", ".join(tags)
    if detail_desc and detail_desc != "정보 없음": return detail_desc[:100] + ("..." if len(detail_desc) > 100 else "") # 길이 제한
    return "특이사항 없음"

# 이미지 관련 함수 (자리만 마련)
def prepare_image_url(value, data_dict=None, json_path=None):
    if value and isinstance(value, str): return IMAGE_BASE_URL + value
    return None

def prepare_image_url_if_type14(value, data_dict=None, json_path=None):
    # grandPlanList[0].imageSrc 를 받았다고 가정. 실제로는 type 확인 로직 필요.
    # 여기서는 data_dict (articleDetail) 와 json_path (grandPlanList[0].imageSrc)를 통해 grandPlanList를 가져와야 함.
    # 또는, get_nested_value에서 grandPlanList[0] 객체 전체를 가져오도록 수정하고, 여기서 type을 확인.
    # 지금은 단순 URL 조합만.
    grand_plan_list = get_nested_value(data_dict, 'articleDetail.grandPlanList', [])
    if grand_plan_list and isinstance(grand_plan_list, list) and len(grand_plan_list) > 0:
        first_plan = grand_plan_list[0]
        if isinstance(first_plan, dict) and first_plan.get('imageType') == "14" and first_plan.get('imageSrc'):
            return IMAGE_BASE_URL + first_plan.get('imageSrc')
    return None


def generate_map_image_placeholder(values, data_dict=None, json_path=None):
    # values는 [latitude, longitude]
    lat = values[0] if len(values) > 0 and values[0] else None
    lon = values[1] if len(values) > 1 and values[1] else None
    if lat and lon: return f"지도 이미지 생성 필요 (위도: {lat}, 경도: {lon})"
    return "위경도 정보 없음"


helper_functions = {
    'format_date_yyyy_mm': format_date_yyyy_mm,
    'get_property_order': get_property_order,
    'format_property_title_text': format_property_title_text,
    'format_household_count': format_household_count,
    'format_floor_count': format_floor_count,
    'format_heating_info': format_heating_info,
    'format_dong_ho_for_ppt': format_dong_ho_for_ppt,
    'format_supply_area_text': format_supply_area_text,
    'format_exclusive_area_text': format_exclusive_area_text,
    'format_room_bathroom_text': format_room_bathroom_text,
    'format_management_fee_for_ppt': format_management_fee_for_ppt,
    'format_price_details_for_ppt': format_price_details_for_ppt,
    'get_first_available_value': get_first_available_value, # 주의: 이 함수는 data_dict와 json_path를 보고 동작해야할 수도 있음
    'format_note_for_ppt': format_note_for_ppt,
    'prepare_image_url': prepare_image_url,
    'prepare_image_url_if_type14': prepare_image_url_if_type14,
    'generate_map_image_placeholder': generate_map_image_placeholder,
}

# --- 4. PPTX 요소 분석 및 수정 함수 ---
def find_shape_by_name(slide_or_group, shape_name):
    for shape in slide_or_group.shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found_in_group = find_shape_by_name(shape, shape_name)
            if found_in_group:
                return found_in_group
    return None


def replace_text_in_shape(shape, placeholder_text_or_none, new_text, font_details=None):
    if not shape or not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    new_text_str = str(new_text if new_text is not None else "")

    text_frame.clear() 
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = new_text_str
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]
        # 필요시 color, italic 등 추가


def add_image_to_shape(shape, image_url_or_path):
    if not shape: return
    print(f"    '{shape.name}'에 이미지 채우기 시도: {image_url_or_path}")
    # 실제 이미지 로드 로직 (주석 해제 필요)
    # try:
    #     if image_url_or_path.startswith("http"):
    #         # import requests
    #         # from io import BytesIO
    #         response = requests.get(image_url_or_path, stream=True)
    #         response.raise_for_status()
    #         image_stream = BytesIO(response.content) # .raw.read() 대신 .content 사용
    #         shape.fill.background() 
    #         shape.fill.solid() # To make sure it's not transparent
    #         shape.fill.fore_color.rgb = RGBColor(255,255,255) # Optional: set background if needed
    #         # For AutoShape, use add_picture, not fill.picture
    #         if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE: # 1 for AUTO_SHAPE
    #              # Ensure the shape is treated as a placeholder for a picture if possible
    #              # This might require the shape to be a picture placeholder in the template
    #              # As a fallback, we try to insert it, but aspect ratio and fit might be off
    #              left, top, width, height = shape.left, shape.top, shape.width, shape.height
    #              shape.element.getparent().remove(shape.element) # Remove old shape
    #              # Add picture in the same location. Note: this changes the shape type.
    #              slide_part = shape.part.package.parts.get_by_name(shape.part.partname)
    #              # This is getting complex, direct picture fill for AutoShape is not straightforward
    #              # Consider making these PICTURE placeholders in the template
    #              # For now, we'll skip direct filling of AutoShape to avoid errors.
    #              print(f"      경고: '{shape.name}' (AUTO_SHAPE)에 대한 자동 이미지 채우기는 현재 지원되지 않습니다. 템플릿에서 그림 자리 표시자로 만드세요.")
    #              # slide_part.add_picture(image_stream, left, top, width, height)
    #         elif hasattr(shape.fill, 'picture'):
    #             shape.fill.picture(image_stream)
    #         print(f"      '{shape.name}'에 URL 이미지 채우기 성공 (가정).")
    #     elif image_url_or_path: # 로컬 경로
    #         # As above, for AutoShape this is problematic
    #         if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
    #             print(f"      경고: '{shape.name}' (AUTO_SHAPE)에 대한 자동 이미지 채우기는 현재 지원되지 않습니다.")
    #         elif hasattr(shape.fill, 'picture'):
    #             shape.fill.background()
    #             shape.fill.picture(image_url_or_path)
    #             print(f"      '{shape.name}'에 로컬 이미지 채우기 성공 (가정).")
    # except Exception as e:
    #     print(f"      오류: '{shape.name}'에 이미지 채우기 중 오류: {e}")
    # else:
    #     print(f"    이미지 경로/URL이 없습니다: {image_url_or_path}")


def analyze_slide_shapes_for_debug(slide_number, slide):
    print(f"\n--- 슬라이드 {slide_number + 1} 요소 분석 (디버깅용) ---")
    if not slide.shapes:
        print("  >> 이 슬라이드에는 도형이 없습니다.")
        return
    
    q = list(slide.shapes) # Use a queue for breadth-first iteration of groups
    
    idx = 0
    while q:
        shape = q.pop(0)
        idx += 1
        shape_name = shape.name if shape.name else "(이름 없음)"
        
        shape_type_val = None
        if hasattr(shape, 'shape_type'): 
            shape_type_val = shape.shape_type
        
        shape_type_name = "알 수 없음"
        if shape_type_val is not None:
            try: 
                shape_type_name = MSO_SHAPE_TYPE(shape_type_val).name 
            except ValueError: 
                shape_type_name = f"커스텀/알수없는 유형 ({shape_type_val})"
            
        text_content = ""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text_content += run.text
        
        print(f"  도형 #{idx}: 이름='{shape_name}', 유형='{shape_type_name} ({shape_type_val})', 텍스트(축약)='{text_content[:50]}...'")
        
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    이 도형은 그룹입니다. 하위 도형:")
            for s_in_group in shape.shapes: # Add shapes in group to the queue
                q.append(s_in_group)
        elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            if shape.has_table:
                table = shape.table
                print(f"      이 도형은 표(Table)입니다: {len(table.rows)}행 x {len(table.columns)}열")
                # for r_idx, row in enumerate(table.rows):
                #     for c_idx, cell in enumerate(row.cells):
                #         print(f"        Cell({r_idx},{c_idx}): '{cell.text[:30]}...'")


# 테이블 셀 텍스트 및 폰트 설정 유틸리티 함수
def set_cell_text(cell, text, font_details=None):
    text_frame = cell.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = str(text if text is not None else "")
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]
        # if font_details.get("color"): font.color.rgb = font_details["color"] # 예: RGBColor(0,0,0)

# 테이블 채우기 함수
def fill_table_from_mappings(table_shape, table_mapping_rules, json_data, client_data):
    if not table_shape or not table_shape.has_table:
        print(f"  경고: '{table_shape.name}'은(는) 표가 아니거나 찾을 수 없습니다.")
        return
    
    table = table_shape.table
    
    for (row_idx, col_idx), (header_text_or_none, json_path_or_list, func_name, font_details) in table_mapping_rules.items():
        if row_idx >= len(table.rows) or col_idx >= len(table.columns):
            print(f"  경고: '{table_shape.name}' 표에 ({row_idx},{col_idx}) 셀이 없습니다.")
            continue
        
        cell = table.cell(row_idx, col_idx)
        final_text = ""
        
        if json_path_or_list == "TABLE_HEADER": # 헤더 텍스트 직접 사용
            final_text = header_text_or_none
        else:
            # 값 가져오기 (json_data 또는 client_data에서)
            # 참고사항/비고 처럼 여러 소스에서 값을 가져와야 하는 경우 get_first_available_value 등에서 처리
            # get_nested_value는 기본적으로 article_json_data (full_json_data)를 대상으로 함
            # client_data의 값을 사용해야 하면, json_path_or_list에 client_data의 키를 직접 사용하고,
            # 아래 get_nested_value 호출 시 data_dict를 client_data로 전달해야 함.
            # 또는, 가공함수에서 data_dict를 client_data로 받도록 설계.
            # 여기서는 json_path에 client_data 키가 올 경우를 대비해 data_source를 선택.

            current_data_source = json_data # 기본값은 full_json_data
            actual_json_path = json_path_or_list

            # 비고: client_data에서 직접 키로 가져오는 경우는 mappings에서 json_path에 해당 키를 직접 명시해야 함.
            # 예를 들어 "참고사항_입력"이 json_path로 오면 client_data에서 찾아야함.
            # 이 부분은 get_nested_value가 client_data도 검색할 수 있도록 하거나,
            # 가공 함수 (예: get_first_available_value) 내에서 client_data를 참조하도록 해야함.
            # 지금은 가공함수에서 여러 값을 받아 처리하는 것으로 가정하고, get_nested_value는 json_data만 봄.
            
            raw_value = None
            if isinstance(actual_json_path, list):
                # 여러 경로에서 값을 가져와 리스트로 만듦 (가공 함수가 리스트를 기대할 때)
                raw_values_list = []
                for jp in actual_json_path:
                    if jp in client_data: # client_data 키 직접 참조
                         raw_values_list.append(client_data.get(jp))
                    else: # json_data (article_json_data) 에서 경로로 탐색
                         raw_values_list.append(get_nested_value(json_data, jp))
                raw_value = raw_values_list
            elif actual_json_path : # 단일 경로
                 if actual_json_path in client_data:
                     raw_value = client_data.get(actual_json_path)
                 else:
                     raw_value = get_nested_value(json_data, actual_json_path)


            final_text = raw_value
            if func_name and func_name in helper_functions:
                try:
                    # 함수 시그니처에 따라 인자 전달 방식 조절 필요
                    # 대부분의 가공함수는 (value, data_dict, json_path)를 받을 수 있도록 설계하는 것이 유연함.
                    # data_dict는 주로 main json_data 이지만, client_data가 필요할 수도 있음.
                    # json_path는 현재 처리중인 경로를 전달.
                    final_text = helper_functions[func_name](raw_value, json_data, actual_json_path)
                except Exception as e:
                    print(f"    오류: 함수 '{func_name}' 실행 중 오류 ({row_idx},{col_idx}): {e}")
                    final_text = "오류"
            
            if isinstance(final_text, list) and not func_name: # 가공 함수 없이 리스트면 join (거의 없을 케이스)
                final_text = ", ".join(map(str, filter(None, final_text)))
            elif final_text is None:
                final_text = ""


        set_cell_text(cell, str(final_text), font_details)
    print(f"  '{table_shape.name}' 표 채우기 완료.")


# --- 슬라이드 복제 함수 (reference 샘플 참조하여 개선) ---
def clone_slide(prs, slide_index):
    """지정된 인덱스의 슬라이드를 복제하여 프레젠테이션에 추가합니다.
    
    reference 파일의 deepcopy 방식을 활용해 슬라이드 디자인을 완벽히 복제합니다.
    """
    from copy import deepcopy
    import os
    
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 이미지 복사를 위한 정보 수집
    image_shapes = {}
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 해당 이미지의 이름과 위치, 크기 정보 저장
            image_shapes[shape.name] = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            # 이미지 파일 경로 추출 시도
            if hasattr(shape, 'image') and hasattr(shape.image, 'filename'):
                image_shapes[shape.name]['filename'] = shape.image.filename
            else:
                image_shapes[shape.name]['filename'] = None
    
    # XML 요소를 완전히 복사하여 슬라이드 디자인 유지
    try:
        for shape in source_slide.shapes:
            el = shape.element
            newel = deepcopy(el)
            # p:extLst 요소가 없을 수 있으므로 예외 처리
            try:
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            except:
                # extLst가 없으면 그냥 마지막에 추가
                new_slide.shapes._spTree.append(newel)
        
        print(f"        슬라이드 {slide_index+1}번 복제 완료. 전체 슬라이드: {len(prs.slides)}개")
    except Exception as e:
        print(f"        슬라이드 복제 중 오류 발생: {e}")
        print("        오류로 인해 기본 슬라이드 복제 방식으로 대체합니다.")
    
    # placeholder 요소 삭제 (Title 1, Subtitle 2 등)
    shapes_to_remove = []
    for shape in new_slide.shapes:
        # placeholder 이거나 Title, Subtitle 등이 이름에 포함된 경우 삭제 대상으로 추가
        if shape.name.startswith('Title') or shape.name.startswith('Subtitle') or \
           shape.name.startswith('Text Placeholder') or 'Placeholder' in shape.name:
            shapes_to_remove.append(shape)
    
    # 삭제 대상 요소 제거
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"        플레이스홀더 삭제: '{shape.name}'")
        except Exception as e:
            print(f"        플레이스홀더 삭제 오류: {e}")
    
    # 특별히 필요한 이미지들을 수동으로 처리
    # 이미지는 직접 파일을 참조해야 합니다. xml 요소만 복사해서는 이미지 파일 자체가 복사되지 않습니다.
    for shape_name in ["img_company_logo"]:
        if shape_name in image_shapes:
            try:
                # 기존 로고 이미지 요소 찾아서 제거 (엑스박스 제거)
                old_logo_shape = None
                for shape in new_slide.shapes:
                    if shape.name == shape_name:
                        old_logo_shape = shape
                        break
                
                if old_logo_shape:
                    # 기존 로고 이미지 요소 제거
                    sp = old_logo_shape._element
                    sp.getparent().remove(sp)
                    print(f"        기존 로고 이미지 제거 완료 ({shape_name})")
                
                # 이미지 파일을 수동으로 추가
                if shape_name == "img_company_logo":
                    # 사용자가 추가한 로고 이미지 파일 사용
                    logo_path = 'NEWFIRST_SYMBOL.png'
                    if os.path.exists(logo_path):
                        from PIL import Image
                        
                        # 원본 이미지 정보
                        shape_info = image_shapes[shape_name]
                        target_width = shape_info['width']
                        target_height = shape_info['height']
                        
                        # 원본 로고 사이즈를 유지하면서 이미지 추가
                        try:
                            # 원본 이미지 정보에서 전달받은 크기를 그대로 사용
                            # 원본 로고의 위치와 크기가 정확히 유지됨
                            new_slide.shapes.add_picture(
                                logo_path,
                                shape_info['left'],
                                shape_info['top'],
                                shape_info['width'],  # 원본 로고의 넓이 그대로 유지
                                shape_info['height']  # 원본 로고의 높이 그대로 유지
                            )
                            print(f"        원본 로고 크기로 이미지 추가 - 넓이: {shape_info['width']}, 높이: {shape_info['height']}")
                            
                            # 영역 정보 출력 (debug용)
                            # print(f"        원본 로고 영역 - left: {shape_info['left']}, top: {shape_info['top']}")
                            
                            
                        except ImportError:
                            # PIL 라이브러리가 없는 경우 기본 방식으로 처리
                            new_slide.shapes.add_picture(
                                logo_path,
                                shape_info['left'],
                                shape_info['top'],
                                target_width,
                                target_height
                            )
                            
                        print(f"        회사 로고 이미지 복제 완료: {logo_path}")
                    else:
                        print(f"        회사 로고 파일을 찾을 수 없습니다: {logo_path}")
            except Exception as e:
                print(f"        이미지 '{shape_name}' 복제 오류: {e}")

    
    # 이제 재사용하기 위해 슬라이드의 텍스트 또는 그림 정보 업데이트 가능
    return new_slide

# --- 5. 메인 PPTX 생성 로직 ---
def generate_presentation(client_data, article_json_list, 매핑규칙):
    """
    여러 매물 정보를 포함하는 PPT 생성 기능
    
    Args:
        client_data (dict): 고객 및 문서 관련 데이터
        article_json_list (list): 여러 매물 데이터의 리스트. 리스트가 아닌 단일 매물 데이터도 처리함
        매핑규칙 (dict): PPT 요소와 데이터 매핑 규칙
    """
    prs = Presentation(TEMPLATE_PATH)
    print(f"\n템플릿 파일 '{TEMPLATE_PATH}' 로드 완료.")

    # 슬라이드 1 채우기
    slide1 = prs.slides[0]
    print("\n--- 슬라이드 1 채우기 시작 ---")
    # 슬라이드 1의 TEXT_BOX 및 TABLE 유형 도형 처리
    for shape_name, mapping_value in 매핑규칙.items(): # value를 먼저 받고, 타입 체크 후 unpack
        if shape_name not in ["txt_document_title", "txt_client_name", "txt_company_name"]:
            continue

        # 매핑 값의 타입이 튜플이고, 길이가 4인지 확인
        if not (isinstance(mapping_value, tuple) and len(mapping_value) == 4):
            print(f"  정보: 슬라이드 1 처리 중 '{shape_name}'의 매핑 값 형식이 올바르지 않아 건너뜁니다. 값: {mapping_value}")
            continue
        
        placeholder_or_cell_coord, json_key, func_name, font_details = mapping_value

        shape_to_fill = find_shape_by_name(slide1, shape_name)
        if not shape_to_fill:
            print(f"  경고: 슬라이드 1에서 '{shape_name}'을(를) 찾지 못했습니다.")
            continue

        value_to_insert = client_data.get(json_key, f"키 '{json_key}' 없음")
        # (슬라이드 1에 가공함수 필요시 여기에 추가)

        if shape_name == "txt_client_name": # 테이블인 경우
            if shape_to_fill.has_table:
                table = shape_to_fill.table
                if table.cell(placeholder_or_cell_coord[0], placeholder_or_cell_coord[1]): # (0,0)
                    set_cell_text(table.cell(placeholder_or_cell_coord[0], placeholder_or_cell_coord[1]), value_to_insert, font_details)
                    print(f"  '{shape_name}' 표의 {placeholder_or_cell_coord} 셀에 '{value_to_insert}' 삽입 및 폰트 적용 완료.")
            else: print(f"  경고: '{shape_name}'은 표가 아닙니다.")
        elif shape_to_fill.has_text_frame: # 텍스트 박스인 경우
            replace_text_in_shape(shape_to_fill, placeholder_or_cell_coord, value_to_insert, font_details)
            print(f"  '{shape_name}' (자리표시자/텍스트박스)에 '{value_to_insert}' 삽입 및 폰트 적용 완료.")
        else:
            print(f"  경고: '{shape_name}'은 처리할 수 없는 유형입니다.")
            
    # 슬라이드 2 채우기
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]
        print("\n--- 슬라이드 2 텍스트 및 이미지 채우기 시작 ---")

        # 1. 매물 타이틀 채우기 (txt_property_title 테이블)
        title_shape_name = "txt_property_title"
        if title_shape_name in 매핑규칙:
            title_coord, title_json_paths, title_func, title_font = 매핑규칙[title_shape_name]
            title_table_shape = find_shape_by_name(slide2, title_shape_name)
            if title_table_shape and title_table_shape.has_table:
                title_values = [get_nested_value(article_json_list[0], jp) for jp in title_json_paths]
                title_text = helper_functions[title_func](title_values) if title_func in helper_functions else ", ".join(map(str,title_values))
                set_cell_text(title_table_shape.table.cell(title_coord[0], title_coord[1]), title_text, title_font)
                print(f"  '{title_shape_name}' 표의 {title_coord} 셀에 '{title_text}' 삽입 및 폰트 적용 완료.")
            else:
                print(f"  경고: 슬라이드 2에서 '{title_shape_name}' 표를 찾지 못했습니다.")

        # 2. 일반 텍스트 박스 채우기 (txt_property_page 등)
        page_shape_name = "txt_property_page"
        if page_shape_name in 매핑규칙:
            page_placeholder, page_json_path, page_func, page_font = 매핑규칙[page_shape_name]
            page_shape = find_shape_by_name(slide2, page_shape_name)
            if page_shape and page_shape.has_text_frame:
                page_raw_val = get_nested_value(article_json_list[0], page_json_path) # json_path가 None일 수 있음
                page_text = helper_functions[page_func](page_raw_val) if page_func in helper_functions else str(page_raw_val or "")
                replace_text_in_shape(page_shape, page_placeholder, page_text, page_font)
                print(f"  '{page_shape_name}' 텍스트 박스에 '{page_text}' 삽입 및 폰트 적용 완료.")
            else:
                print(f"  경고: 슬라이드 2에서 '{page_shape_name}' 텍스트 상자를 찾지 못했습니다.")

        # 3. 주요 정보 테이블들 채우기
        table_names_to_fill = ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]
        for table_name in table_names_to_fill:
            if table_name in 매핑규칙:
                table_mapping_rules = 매핑규칙[table_name] # 이것은 셀별 매핑 딕셔너리
                table_shape_obj = find_shape_by_name(slide2, table_name)
                if table_shape_obj:
                    fill_table_from_mappings(table_shape_obj, table_mapping_rules, article_json_list[0], client_data)
                else:
                    print(f"  경고: 슬라이드 2에서 '{table_name}' 표를 찾지 못했습니다.")
            else:
                print(f"  정보: '{table_name}'에 대한 매핑 규칙이 없습니다.")
        
        # 4. 이미지 채우기 (주석 처리된 상태)
        print("\n--- 슬라이드 2 이미지 채우기 (실제 로드는 주석 처리됨) ---")
        image_shapes_to_fill = ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]
        for img_shape_name in image_shapes_to_fill:
            if img_shape_name in 매핑규칙:
                _, img_json_path, img_func_name, _ = 매핑규칙[img_shape_name]
                img_shape = find_shape_by_name(slide2, img_shape_name)
                if img_shape:
                    raw_img_value = None
                    if isinstance(img_json_path, list): # 위경도 같은 경우
                        raw_img_value = [get_nested_value(article_json_list[0], jp) for jp in img_json_path]
                    else:
                        raw_img_value = get_nested_value(article_json_list[0], img_json_path)
                    
                    image_url_or_data = raw_img_value
                    if img_func_name and img_func_name in helper_functions:
                        image_url_or_data = helper_functions[img_func_name](raw_img_value, article_json_list[0], img_json_path)
                    
                    add_image_to_shape(img_shape, image_url_or_data)
                else:
                    print(f"  경고: 슬라이드 2에서 이미지 도형 '{img_shape_name}'을(를) 찾지 못했습니다.")

    # 여러 매물 데이터 처리
    # 리스트가 아닌 단일 매물 데이터도 리스트로 변환하여 처리
    if not isinstance(article_json_list, list):
        article_json_list = [article_json_list]
    
    # 추가 매물 정보 슬라이드 생성
    for idx, article_data in enumerate(article_json_list[1:], start=1):
        # 슬라이드 2를 복제하여 새 슬라이드 생성
        new_slide = clone_slide(prs, 1)  # 슬라이드 인덱스 1을 복제 (0부터 시작하므로 인덱스 1은 두 번째 슬라이드임)

        # 슬라이드 번호 변경 (매물 순번)
        prop_title_shape = find_shape_by_name(new_slide, "txt_property_title")
        if prop_title_shape and hasattr(prop_title_shape, 'has_table') and prop_title_shape.has_table:
            # 매물 번호 증가
            div_name = get_nested_value(article_data, "articleDetail.divisionName", "")
            apt_name = get_nested_value(article_data, "articleDetail.aptName", "")
            new_title = f"No.{idx+1} [{div_name}] {apt_name}"
            set_cell_text(prop_title_shape.table.cell(0, 0), new_title, {"font_name": "나눔고딕", "font_size": Pt(18)})
        
        # 매물 페이지 번호 업데이트
        page_shape = find_shape_by_name(new_slide, "txt_property_page")
        if page_shape and hasattr(page_shape, 'has_text_frame') and page_shape.has_text_frame:
            replace_text_in_shape(page_shape, None, str(idx+1), {"font_name": "나눔고딕", "font_size": Pt(10)})
        
        # 매물 데이터 테이블 채우기
        # --- 단지 정보 테이블 채우기 ---
        tables_to_fill = ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]
        for table_name in tables_to_fill:
            if table_name in 매핑규칙:
                table_mapping_rules = 매핑규칙[table_name]
                table_shape_obj = find_shape_by_name(new_slide, table_name)
                if table_shape_obj and hasattr(table_shape_obj, 'has_table') and table_shape_obj.has_table:
                    fill_table_from_mappings(table_shape_obj, table_mapping_rules, article_data, client_data)
        
        # --- 이미지 도형 채우기 ---
        image_shapes_to_fill = ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]
        for img_shape_name in image_shapes_to_fill:
            if img_shape_name in 매핑규칙:
                _, img_json_path, img_func_name, _ = 매핑규칙[img_shape_name]
                img_shape = find_shape_by_name(new_slide, img_shape_name)
                
                # 이미지 로드 오류 방지
                try:
                    if img_shape:
                        if isinstance(img_json_path, list):
                            raw_img_value = [get_nested_value(article_data, jp) for jp in img_json_path]
                        else:
                            raw_img_value = get_nested_value(article_data, img_json_path)
                        
                        image_url_or_data = raw_img_value
                        if img_func_name and img_func_name in helper_functions:
                            image_url_or_data = helper_functions[img_func_name](raw_img_value, article_data, img_json_path)
                        
                        add_image_to_shape(img_shape, image_url_or_data)
                except Exception as e:
                    print(f"        이미지 처리 오류 ({img_shape_name}): {e} - 무시하고 계속")
    
    # 결과물 반환
    prs.save(OUTPUT_PATH_DEBUG)
    print(f"\n최종 PPTX 파일 '{OUTPUT_PATH_DEBUG}' 저장 완료. 총 {len(article_json_list)+1}개 슬라이드 생성됨.")


if __name__ == "__main__":
    try:
        prs_debug = Presentation(TEMPLATE_PATH)
        if len(prs_debug.slides) > 0:
            analyze_slide_shapes_for_debug(0, prs_debug.slides[0])
        if len(prs_debug.slides) > 1:
            analyze_slide_shapes_for_debug(1, prs_debug.slides[1])
    except Exception as e:
        print(f"디버깅 정보 출력 중 오류: {e}")
    
    print("\n--- 메인 생성 로직 시작 ---")
    
    # 여러 매물 데이터를 리스트로 묶어서 전달
    article_data_list = [first_property_data, second_property_data, third_property_data]
    
    print(f"총 {len(article_data_list)}개의 매물 데이터로 PPT 생성을 시작합니다.")
    generate_presentation(client_and_document_data, article_data_list, mappings)