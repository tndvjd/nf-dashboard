from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
import datetime
import math
from io import BytesIO
# import requests # 실제 이미지 로드 시 주석 해제
# from io import BytesIO # 실제 이미지 로드 시 주석 해제

TEMPLATE_PATH = "ppt_template.pptx" # FastAPI에서는 이 경로를 어떻게 처리할지 고려 필요 (예: 설정 파일, 환경 변수)
IMAGE_BASE_URL = "https://land.naver.com"

# --- 매핑 규칙 정의 ---
mappings = {
    # 슬라이드 1
    "txt_document_title": ("{{document_title}}", "문서명", None, {"font_name": "나눔고딕", "font_size": Pt(24)}),
    "txt_client_name": ((0,0), "고객명", None, {"font_name": "나눔고딕", "font_size": Pt(18)}), # 테이블 셀
    "txt_company_name": ("{{company_name}}", "회사명", None, {"font_name": "나눔고딕", "font_size": Pt(16)}),

    # 슬라이드 2
    "txt_property_title": ((0,0), ["articleDetail.divisionName", "articleDetail.aptName"], 'format_property_title_text', {"font_name": "나눔고딕", "font_size": Pt(18)}),
    "txt_property_page": ("{{매물순번_표시}}", None, 'get_property_order', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    "tbl_complex_info": {
        (0, 0): ("단지주소", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, "articleDetail.exposureAddress", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("준공연도", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleDetail.aptUseApproveYmd", 'format_date_yyyy_mm', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("총세대수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.aptHouseholdCount", 'format_household_count', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("총층수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, "articleFloor.buildingHighestFloor", 'format_floor_count', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("난방방식", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.aptHeatMethodTypeName", "articleDetail.aptHeatFuelTypeName"], 'format_heating_info', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "tbl_property_detail_1": {
        (0, 0): ("동·호수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articleDetail.buildingName", "articleAddition.floorInfo"], 'format_dong_ho_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("계약면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleSpace.supplySpace", 'format_supply_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("전용면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleSpace.exclusiveSpace", 'format_exclusive_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("방수/욕실수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.roomCount", "articleDetail.bathroomCount"], 'format_room_bathroom_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("방향", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, "articleAddition.direction", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "tbl_property_detail_2": {
        (0, 0): ("보증금/월세", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articlePrice.warrantPrice", "articlePrice.rentPrice", "articleAddition.dealOrWarrantPrc", "articleAddition.rentPrc"], 'format_price_details_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("기본관리비", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "administrationCostInfo.chargeCodeType", 'format_management_fee_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("입주가능일", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.moveInTypeName", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("참고사항", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.articleFeatureDescription", "참고사항_입력"], 'get_first_available_value', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("비고", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.tagList", "articleDetail.detailDescription", "비고_입력"], 'format_note_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "img_complex_view": (None, "articleAddition.representativeImgUrl", 'prepare_image_url', None),
    "img_complex_floorplan": (None, "articleDetail.grandPlanList[0].imageSrc", 'prepare_image_url_if_type14', None),
    "img_complex_mapimg": (None, ["articleDetail.latitude", "articleDetail.longitude"], 'generate_map_image_placeholder', None),
}

# --- 데이터 가공 헬퍼 함수들 ---
def get_nested_value(data_dict, keys_string_or_list, default=""):
    if keys_string_or_list is None: return default
    if isinstance(keys_string_or_list, list):
        for keys_string in keys_string_or_list:
            val = get_single_nested_value(data_dict, keys_string, None)
            if val is not None and val != "": return val
        return default
    return get_single_nested_value(data_dict, keys_string_or_list, default)

def get_single_nested_value(data_dict, keys_string, default=""):
    if keys_string is None:
        return default
    keys = keys_string.split('.')
    value = data_dict
    try:
        for key_part in keys:
            if '[' in key_part and key_part.endswith(']'):
                key = key_part.split('[')[0]
                index = int(key_part.split('[')[1][:-1])
                if key in value and isinstance(value[key], list) and len(value[key]) > index:
                    value = value[key][index]
                else: return default
            elif isinstance(value, dict) and key_part in value:
                value = value[key_part]
            else: return default
        return value if value is not None else default
    except (KeyError, IndexError, TypeError, AttributeError): return default

def format_date_yyyy_mm(value, data_dict=None, json_path=None):
    if not value or len(str(value)) != 8: return ""
    try:
        dt_obj = datetime.datetime.strptime(str(value), "%Y%m%d")
        return f"{dt_obj.year}년 {dt_obj.month}월"
    except ValueError: return ""

def calculate_pyeong(m2_value, precision=1):
    if isinstance(m2_value, (int, float)) and m2_value > 0:
        return round(m2_value / 3.3058, precision)
    return ""

def format_property_title_text(values, data_dict=None, json_path=None):
    prop_order = "1"
    div_name = values[0] if len(values) > 0 and values[0] else ""
    apt_name = values[1] if len(values) > 1 and values[1] else ""
    return f"No.{prop_order} [{div_name}] {apt_name}"

def get_property_order(value=None, data_dict=None, json_path=None):
    return "1"

def format_household_count(value, data_dict=None, json_path=None):
    return f"{value}세대" if value else ""

def format_floor_count(value, data_dict=None, json_path=None):
    return f"{value}층" if value else ""

def format_heating_info(values, data_dict=None, json_path=None):
    method = values[0] if len(values) > 0 and values[0] else ""
    fuel = values[1] if len(values) > 1 and values[1] else ""
    if method and fuel: return f"{method}, {fuel}"
    return method or fuel or ""

def format_dong_ho_for_ppt(values, data_dict=None, json_path=None):
    building_name = values[0] if len(values) > 0 and values[0] else ""
    floor_info_raw = values[1] if len(values) > 1 and values[1] else ""
    
    floor_display = ""
    if isinstance(floor_info_raw, str):
        if "중층" in floor_info_raw or "고층" in floor_info_raw or "저층" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0]
        elif "/" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0] + "층"
        else: floor_display = floor_info_raw
    else: floor_display = str(floor_info_raw)
    
    return f"{building_name} {floor_display}".strip()

def format_supply_area_text(supply_m2_val, data_dict=None, json_path=None):
    supply_py_val = calculate_pyeong(supply_m2_val)
    return f"{supply_m2_val}㎡ ({supply_py_val}평)" if supply_m2_val else ""

def format_exclusive_area_text(exclusive_m2_val, data_dict=None, json_path=None):
    exclusive_py_val = calculate_pyeong(exclusive_m2_val)
    return f"{exclusive_m2_val}㎡ ({exclusive_py_val}평)" if exclusive_m2_val else ""

def format_room_bathroom_text(values, data_dict=None, json_path=None):
    room = values[0] if len(values) > 0 and values[0] else "0"
    bath = values[1] if len(values) > 1 and values[1] else "0"
    return f"{room} / {bath}"

def format_management_fee_for_ppt(value, data_dict=None, json_path=None):
    return "확인 어려움"

def format_price_details_for_ppt(values, data_dict=None, json_path=None):
    w_price_num = values[0] if len(values) > 0 and values[0] else 0
    r_price_num = values[1] if len(values) > 1 and values[1] else 0

    formatted_w = ""
    if isinstance(w_price_num, (int, float)) and w_price_num > 0:
        if w_price_num >= 10000:
            eok = int(w_price_num // 10000)
            manwon_part = int(w_price_num % 10000)
            formatted_w = f"{eok}억"
            if manwon_part > 0: formatted_w += f" {manwon_part:,}"
        else: formatted_w = f"{int(w_price_num):,}"
    
    formatted_r = ""
    if isinstance(r_price_num, (int, float)) and r_price_num > 0:
        formatted_r = f"{int(r_price_num):,}"
        
    if formatted_w and formatted_r: return f"{formatted_w} / {formatted_r} (만원)"
    if formatted_w: return f"{formatted_w} (만원)"
    if formatted_r: return f" / {formatted_r} (만원)"
    return "가격 정보 없음"

def get_first_available_value(values, data_dict=None, json_path=None):
    for val_candidate_key in values:
        val = val_candidate_key
        if val and val != "정보 없음": return str(val)
    return "내용 없음"

def format_note_for_ppt(values, data_dict=None, json_path=None):
    tags = values[0] if len(values) > 0 and isinstance(values[0], list) else []
    detail_desc = values[1] if len(values) > 1 and values[1] else ""
    manual_note = values[2] if len(values) > 2 and values[2] else ""

    if manual_note and manual_note != "정보 없음": return manual_note
    if tags: return ", ".join(tags)
    if detail_desc and detail_desc != "정보 없음": return detail_desc[:100] + ("..." if len(detail_desc) > 100 else "")
    return "특이사항 없음"

def prepare_image_url(value, data_dict=None, json_path=None):
    if value and isinstance(value, str): return IMAGE_BASE_URL + value
    return None

def prepare_image_url_if_type14(value, data_dict=None, json_path=None):
    grand_plan_list = get_nested_value(data_dict, 'articleDetail.grandPlanList', [])
    if grand_plan_list and isinstance(grand_plan_list, list) and len(grand_plan_list) > 0:
        first_plan = grand_plan_list[0]
        if isinstance(first_plan, dict) and first_plan.get('imageType') == "14" and first_plan.get('imageSrc'):
            return IMAGE_BASE_URL + first_plan.get('imageSrc')
    return None

def generate_map_image_placeholder(values, data_dict=None, json_path=None):
    lat = values[0] if len(values) > 0 and values[0] else None
    lon = values[1] if len(values) > 1 and values[1] else None
    if lat and lon: return f"지도 이미지 생성 필요 (위도: {lat}, 경도: {lon})"
    return "위경도 정보 없음"

helper_functions = {
    'format_date_yyyy_mm': format_date_yyyy_mm,
    'get_property_order': get_property_order,
    'format_property_title_text': format_property_title_text,
    'format_household_count': format_household_count,
    'format_floor_count': format_floor_count,
    'format_heating_info': format_heating_info,
    'format_dong_ho_for_ppt': format_dong_ho_for_ppt,
    'format_supply_area_text': format_supply_area_text,
    'format_exclusive_area_text': format_exclusive_area_text,
    'format_room_bathroom_text': format_room_bathroom_text,
    'format_management_fee_for_ppt': format_management_fee_for_ppt,
    'format_price_details_for_ppt': format_price_details_for_ppt,
    'get_first_available_value': get_first_available_value,
    'format_note_for_ppt': format_note_for_ppt,
    'prepare_image_url': prepare_image_url,
    'prepare_image_url_if_type14': prepare_image_url_if_type14,
    'generate_map_image_placeholder': generate_map_image_placeholder,
}

# --- PPTX 요소 분석 및 수정 함수 ---
def find_shape_by_name(slide_or_group, shape_name):
    for shape in slide_or_group.shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP: # MSO_SHAPE_TYPE.GROUP은 그룹 도형을 의미합니다.
            found_in_group = find_shape_by_name(shape, shape_name)
            if found_in_group:
                return found_in_group
    return None

def replace_text_in_shape(shape, placeholder_text_or_none, new_text, font_details=None):
    if not shape or not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    new_text_str = str(new_text if new_text is not None else "")

    text_frame.clear() 
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = new_text_str
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]

def add_image_to_shape(shape, image_url_or_path):
    if not shape: return
    # print(f"    '{shape.name}'에 이미지 채우기 시도: {image_url_or_path}")
    # 실제 이미지 로드 로직은 API 서버 환경에 맞게 재구성 필요 (예: 로컬 파일 접근, 외부 URL 접근 권한 등)
    # 현재는 주석 처리된 상태로 유지
    pass

def set_cell_text(cell, text, font_details=None):
    text_frame = cell.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = str(text if text is not None else "")
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]

def fill_table_from_mappings(table_shape, table_mapping_rules, json_data, client_data):
    if not table_shape or not table_shape.has_table:
        # print(f"  경고: '{table_shape.name}'은(는) 표가 아니거나 찾을 수 없습니다.")
        return
    
    table = table_shape.table
    
    for (row_idx, col_idx), (header_text_or_none, json_path_or_list, func_name, font_details) in table_mapping_rules.items():
        if row_idx >= len(table.rows) or col_idx >= len(table.columns):
            # print(f"  경고: '{table_shape.name}' 표에 ({row_idx},{col_idx}) 셀이 없습니다.")
            continue
        
        cell = table.cell(row_idx, col_idx)
        final_text = ""
        
        if json_path_or_list == "TABLE_HEADER":
            final_text = header_text_or_none
        else:
            current_data_source = json_data
            actual_json_path = json_path_or_list
            
            raw_value = None
            if isinstance(actual_json_path, list):
                raw_values_list = []
                for jp in actual_json_path:
                    if jp in client_data:
                         raw_values_list.append(client_data.get(jp))
                    else:
                         raw_values_list.append(get_nested_value(json_data, jp))
                raw_value = raw_values_list
            elif actual_json_path :
                 if actual_json_path in client_data:
                     raw_value = client_data.get(actual_json_path)
                 else:
                     raw_value = get_nested_value(json_data, actual_json_path)

            final_text = raw_value
            if func_name and func_name in helper_functions:
                try:
                    final_text = helper_functions[func_name](raw_value, json_data, actual_json_path)
                except Exception as e:
                    # print(f"    오류: 함수 '{func_name}' 실행 중 오류 ({row_idx},{col_idx}): {e}")
                    final_text = "오류"
            
            if isinstance(final_text, list) and not func_name:
                final_text = ", ".join(map(str, filter(None, final_text)))
            elif final_text is None:
                final_text = ""
        
        set_cell_text(cell, str(final_text), font_details)
    # print(f"  '{table_shape.name}' 표 채우기 완료.")

# --- 슬라이드 복제 헬퍼 함수 ---
def clone_slide(prs, slide_index):
    """지정된 인덱스의 슬라이드를 복제하여 프레젠테이션에 추가합니다."""
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 모든 도형 복제
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # 테이블 복제
            old_table = shape.table
            new_table = new_slide.shapes.add_table(
                old_table.rows.count, 
                old_table.columns.count,
                shape.left, shape.top, shape.width, shape.height
            ).table
            
            # 셀 스타일 및 내용 복제
            for i, row in enumerate(old_table.rows):
                for j, cell in enumerate(row.cells):
                    new_cell = new_table.cell(i, j)
                    new_cell.fill.solid()
                    new_cell.fill.fore_color.rgb = cell.fill.fore_color.rgb
                    
                    # 텍스트 프레임의 속성(텍스트 박스) 복제
                    for paragraph in cell.text_frame.paragraphs:
                        new_paragraph = new_cell.text_frame.paragraphs[0] if i == 0 and j == 0 else new_cell.text_frame.add_paragraph()
                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run()
                            new_run.text = run.text
                            
                            # 서식 복제
                            if hasattr(run.font, 'bold'): new_run.font.bold = run.font.bold
                            if hasattr(run.font, 'italic'): new_run.font.italic = run.font.italic
                            if hasattr(run.font, 'underline'): new_run.font.underline = run.font.underline
                            if hasattr(run.font, 'size'): new_run.font.size = run.font.size
                            if hasattr(run.font, 'name'): new_run.font.name = run.font.name
                            if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'): 
                                new_run.font.color.rgb = run.font.color.rgb
            
            # 테이블 속성 복제
            for i in range(len(old_table.columns)):
                new_table.columns[i].width = old_table.columns[i].width
            for i in range(len(old_table.rows)):
                new_table.rows[i].height = old_table.rows[i].height
            
            # 테이블 이름 유지
            for shape_name in ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]:
                if shape.name == shape_name:
                    new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                    break
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 이미지 복제
            image_path = shape.image.filename if hasattr(shape.image, 'filename') else None
            if image_path:
                pic = new_slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                # 이미지 이름 유지
                for shape_name in ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:pic[last()]')[0].set('name', shape_name)
                        break
            else:
                # 이미지 경로를 찾을 수 없는 경우 빈 이미지 자리 유지
                placeholder = new_slide.shapes.add_shape(
                    MSO_SHAPE_TYPE.RECTANGLE, shape.left, shape.top, shape.width, shape.height
                )
                # 이름 유지
                for shape_name in ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                        break
        
        else:
            # 기타 도형 복제 (텍스트 박스 등)
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(left, top, width, height)
                new_text_frame = new_shape.text_frame
                old_text_frame = shape.text_frame
                
                # 텍스트 프레임 속성 복제
                new_text_frame.word_wrap = old_text_frame.word_wrap
                
                # 텍스트 내용 복제
                for i, paragraph in enumerate(old_text_frame.paragraphs):
                    new_paragraph = new_text_frame.paragraphs[0] if i == 0 else new_text_frame.add_paragraph()
                    new_paragraph.alignment = paragraph.alignment
                    
                    for run in paragraph.runs:
                        new_run = new_paragraph.add_run()
                        new_run.text = run.text
                        
                        # 서식 복제
                        if hasattr(run.font, 'bold'): new_run.font.bold = run.font.bold
                        if hasattr(run.font, 'italic'): new_run.font.italic = run.font.italic
                        if hasattr(run.font, 'size'): new_run.font.size = run.font.size
                        if hasattr(run.font, 'name'): new_run.font.name = run.font.name
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'): 
                            new_run.font.color.rgb = run.font.color.rgb
                
                # 텍스트박스 이름 유지
                for shape_name in ["txt_property_title", "txt_property_page"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                        break
    
    # 슬라이드 속성 복제 (배경 등)
    if hasattr(source_slide, 'background') and hasattr(new_slide, 'background'):
        if hasattr(source_slide.background, 'fill') and hasattr(new_slide.background, 'fill'):
            if source_slide.background.fill.type == source_slide.background.fill.SOLID:
                new_slide.background.fill.solid()
                new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
    
    return new_slide

# --- 메인 PPTX 생성 로직 ---
def generate_presentation_logic(client_data: dict, article_json_data: dict, 매핑규칙: dict):
    prs = Presentation(TEMPLATE_PATH)
    # print(f"템플릿 파일 '{TEMPLATE_PATH}' 로드 완료.")

    # 슬라이드 1 (표지) 채우기
    slide1 = prs.slides[0]
    for shape_name, mapping_value in 매핑규칙.items():
        if shape_name not in ["txt_document_title", "txt_client_name", "txt_company_name"]:
            continue

        if not (isinstance(mapping_value, tuple) and len(mapping_value) == 4):
            continue
        
        placeholder_or_cell_coord, json_key, func_name, font_details = mapping_value
        shape_to_fill = find_shape_by_name(slide1, shape_name)
        if not shape_to_fill:
            continue

        value_to_insert = client_data.get(json_key, f"키 '{json_key}' 없음")

        if shape_name == "txt_client_name":
            if shape_to_fill.has_table:
                table = shape_to_fill.table
                if table.cell(placeholder_or_cell_coord[0], placeholder_or_cell_coord[1]):
                    set_cell_text(table.cell(placeholder_or_cell_coord[0], placeholder_or_cell_coord[1]), value_to_insert, font_details)
        elif shape_to_fill.has_text_frame:
            replace_text_in_shape(shape_to_fill, placeholder_or_cell_coord, value_to_insert, font_details)

    # 슬라이드 2 (첫 번째 매물) 채우기
    slide2 = prs.slides[1]
    
    # 매물 데이터 채우기를 위한 함수 정의
    def fill_property_data(slide, article_data, client_data, mapping_rules):
        # --- 슬라이드 텍스트 및 이미지 채우기 ---
        tables_to_fill = ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]
        for table_name in tables_to_fill:
            table_shape = find_shape_by_name(slide, table_name)
            if table_shape and table_shape.has_table and table_name in mapping_rules:
                fill_table_from_mappings(table_shape, mapping_rules[table_name], article_data, client_data)

        # --- 슬라이드 이미지 채우기 (실제 로드는 주석 처리됨) ---
        image_shapes_to_fill = ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]
        for img_shape_name in image_shapes_to_fill:
            if img_shape_name in mapping_rules:
                _, img_json_path, img_func_name, _ = mapping_rules[img_shape_name]
                img_shape = find_shape_by_name(slide, img_shape_name)
                img_func = helper_functions.get(img_func_name)
                if img_shape and img_func:
                    img_url = img_func(get_nested_value(article_data, img_json_path), article_data, img_json_path)
                    if img_url:
                        # 실제 이미지 로드 (주석 처리됨)
                        # print(f"이미지 로드 시도: {img_url}")
                        # add_image_to_shape(img_shape, img_url)
                        pass
    
    # 제목 업데이트: 첫 번째 매물 정보 업데이트 (슬라이드 2)    
    title_shape = find_shape_by_name(slide2, "txt_property_title")
    if title_shape and title_shape.has_table:
        div_name = get_nested_value(article_json_data, "articleDetail.divisionName", "")
        apt_name = get_nested_value(article_json_data, "articleDetail.aptName", "")
        title_text = f"No.1 [{div_name}] {apt_name}"
        cell = title_shape.table.cell(0, 0)
        set_cell_text(cell, title_text, {"font_name": "나눔고딕", "font_size": Pt(18)})
    
    # 페이지 번호 업데이트
    page_shape = find_shape_by_name(slide2, "txt_property_page")
    if page_shape:
        replace_text_in_shape(page_shape, None, "1", {"font_name": "나눔고딕", "font_size": Pt(10)})
        
    # 슬라이드 2 데이터 채우기
    fill_property_data(slide2, article_json_data, client_data, 매핑규칙)

    # 추가 매물 데이터 샘플 (실제로는 API 요청에서 받아올 수 있음)
    additional_article_data = [
        {
            "articleDetail": {
                "divisionName": "아파트",
                "aptName": "래미안 퍼스티지",
                "exposureAddress": "서울특별시 강남구 신사동 123-45",
                "aptUseApproveYmd": "20150312",
                "aptHouseholdCount": "320",
                "aptHeatMethodTypeName": "개별난방",
                "aptHeatFuelTypeName": "도시가스",
                "buildingName": "101동",
                "roomCount": "3",
                "bathroomCount": "2",
                "moveInTypeName": "즉시 입주 가능",
                "articleFeatureDescription": "신사역 도보 5분 거리, 리모델링 완료된 남향 코너 호실",
                "detailDescription": "채광 좋은 거실, 입주 즉시 가능한 깔끔한 아파트입니다.",
                "tagList": ["남향", "역세권", "반려동물 가능", "신규 등록"],
                "grandPlanList": [{
                    "imageSrc": "/타입1_평면도.jpg", 
                    "imageType": "PLAN"
                }],
                "latitude": "37.5167",
                "longitude": "127.0402"
            },
            "articleAddition": {
                "floorInfo": "12/15",
                "direction": "남향",
                "dealOrWarrantPrc": "25000",
                "rentPrc": "120",
                "representativeImgUrl": "/매물사진2.jpg"
            },
            "articleFloor": {
                "totalFloorCount": "15",
                "buildingHighestFloor": "15"
            },
            "articlePrice": {
                "warrantPrice": 25000,
                "rentPrice": 120
            },
            "articleSpace": {
                "supplySpace": 112.32,
                "exclusiveSpace": 84.83
            },
            "articlePhotos": [
                {"imageSrc": "/매물사진2_1.jpg", "imageType": "INTERNAL"},
                {"imageSrc": "/매물사진2_2.jpg", "imageType": "INTERNAL"}
            ],
            "administrationCostInfo": {
                "chargeCodeType": "INCLUDE_ALL"
            },
            "참고사항_입력": "관리비에 수도/전기/가스 포함",
            "비고_입력": "주차 2대 가능"
        },
        {
            "articleDetail": {
                "divisionName": "오피스텔",
                "aptName": "센트럴파크 타워",
                "exposureAddress": "서울특별시 강남구 역삼동 789-10",
                "aptUseApproveYmd": "20180525",
                "aptHouseholdCount": "192",
                "aptHeatMethodTypeName": "중앙난방",
                "aptHeatFuelTypeName": "지역난방",
                "buildingName": "A동",
                "roomCount": "2",
                "bathroomCount": "1",
                "moveInTypeName": "2개월 후 입주 가능",
                "articleFeatureDescription": "역삼역 초역세권, 풀옵션 신축급 오피스텔",
                "detailDescription": "2년 계약 시 중개수수료 할인 가능합니다.",
                "tagList": ["풀옵션", "신축", "주차가능", "24시간 보안"],
                "grandPlanList": [{
                    "imageSrc": "/타입2_평면도.jpg", 
                    "imageType": "PLAN"
                }],
                "latitude": "37.5006",
                "longitude": "127.0359"
            },
            "articleAddition": {
                "floorInfo": "8/18",
                "direction": "동향",
                "dealOrWarrantPrc": "15000",
                "rentPrc": "90",
                "representativeImgUrl": "/매물사진3.jpg"
            },
            "articleFloor": {
                "totalFloorCount": "18",
                "buildingHighestFloor": "18"
            },
            "articlePrice": {
                "warrantPrice": 15000,
                "rentPrice": 90
            },
            "articleSpace": {
                "supplySpace": 63.6,
                "exclusiveSpace": 45.5
            },
            "articlePhotos": [
                {"imageSrc": "/매물사진3_1.jpg", "imageType": "INTERNAL"},
                {"imageSrc": "/매물사진3_2.jpg", "imageType": "INTERNAL"}
            ],
            "administrationCostInfo": {
                "chargeCodeType": "INCLUDE_PART"
            },
            "참고사항_입력": "보안키 및 디지털 도어락 설치",
            "비고_입력": "공용 사우나, 피트니스센터 이용 가능"
        }
    ]
    
    # 추가 매물 정보 슬라이드 생성
    for idx, additional_data in enumerate(additional_article_data, start=1):
        # 슬라이드 2를 복제하여 새 슬라이드 생성
        new_slide = clone_slide(prs, 1)  # 슬라이드 인덱스 1을 복제 (0부터 시작하므로 인덱스 1은 두 번째 슬라이드임)

        # 슬라이드 번호 변경 (매물 순번)
        prop_title_shape = find_shape_by_name(new_slide, "txt_property_title")
        if prop_title_shape and prop_title_shape.has_table:
            # 매물 번호 증가
            div_name = get_nested_value(additional_data, "articleDetail.divisionName", "")
            apt_name = get_nested_value(additional_data, "articleDetail.aptName", "")
            new_title = f"No.{idx+1} [{div_name}] {apt_name}"
            set_cell_text(prop_title_shape.table.cell(0, 0), new_title, {"font_name": "나눔고딕", "font_size": Pt(18)})
        
        # 매물 페이지 번호 업데이트
        page_shape = find_shape_by_name(new_slide, "txt_property_page")
        if page_shape:
            replace_text_in_shape(page_shape, None, str(idx+1), {"font_name": "나눔고딕", "font_size": Pt(10)})
        
        # 복제된 슬라이드에 새로운 추가 데이터 적용
        fill_property_data(new_slide, additional_data, client_data, 매핑규칙)
    
    # 결과물 반환
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer