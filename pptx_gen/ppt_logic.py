from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
import datetime
import math
from io import BytesIO
import copy # Added for deepcopy, though not used in the primary clone_slide logic directly
# import requests # 실제 이미지 로드 시 주석 해제
# from io import BytesIO # 실제 이미지 로드 시 주석 해제

IMAGE_BASE_URL = "https://land.naver.com"

# --- 매핑 규칙 정의 ---
mappings = {
    # 슬라이드 1
    "txt_document_title": ("{{document_title}}", "문서명", None, {"font_name": "나눔고딕", "font_size": Pt(24)}),
    "txt_client_name": ((0,0), "고객명", None, {"font_name": "나눔고딕", "font_size": Pt(18)}), # 테이블 셀
    "txt_company_name": ("{{company_name}}", "회사명", None, {"font_name": "나눔고딕", "font_size": Pt(16)}),

    # 슬라이드 2
    "txt_property_title": ((0,0), ["articleDetail.divisionName", "articleDetail.aptName"], 'format_property_title_text', {"font_name": "나눔고딕", "font_size": Pt(18)}),
    "txt_property_page": ("{{매물순번_표시}}", None, 'get_property_order', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    "tbl_complex_info": {
        (0, 0): ("단지주소", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, "articleDetail.exposureAddress", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("준공연도", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleDetail.aptUseApproveYmd", 'format_date_yyyy_mm', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("총세대수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.aptHouseholdCount", 'format_household_count', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("총층수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, "articleFloor.buildingHighestFloor", 'format_floor_count', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("난방방식", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.aptHeatMethodTypeName", "articleDetail.aptHeatFuelTypeName"], 'format_heating_info', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "tbl_property_detail_1": {
        (0, 0): ("동·호수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articleDetail.buildingName", "articleAddition.floorInfo"], 'format_dong_ho_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("계약면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "articleSpace.supplySpace", 'format_supply_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("전용면적", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleSpace.exclusiveSpace", 'format_exclusive_area_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("방수/욕실수", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.roomCount", "articleDetail.bathroomCount"], 'format_room_bathroom_text', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("방향", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, "articleAddition.direction", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "tbl_property_detail_2": {
        (0, 0): ("보증금/월세", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (0, 1): (None, ["articlePrice.warrantPrice", "articlePrice.rentPrice", "articleAddition.dealOrWarrantPrc", "articleAddition.rentPrc"], 'format_price_details_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (1, 0): ("기본관리비", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (1, 1): (None, "administrationCostInfo.chargeCodeType", 'format_management_fee_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (2, 0): ("입주가능일", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (2, 1): (None, "articleDetail.moveInTypeName", None, {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (3, 0): ("참고사항", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (3, 1): (None, ["articleDetail.articleFeatureDescription", "참고사항_입력"], 'get_first_available_value', {"font_name": "나눔고딕", "font_size": Pt(10)}),
        (4, 0): ("비고", "TABLE_HEADER", None, {"font_name": "나눔고딕", "font_size": Pt(10), "bold": True}),
        (4, 1): (None, ["articleDetail.tagList", "articleDetail.detailDescription", "비고_입력"], 'format_note_for_ppt', {"font_name": "나눔고딕", "font_size": Pt(10)}),
    },
    "img_complex_view": (None, "articleAddition.representativeImgUrl", 'prepare_image_url', None),
    "img_complex_floorplan": (None, "articleDetail.grandPlanList[0].imageSrc", 'prepare_image_url_if_type14', None),
    "img_complex_mapimg": (None, ["articleDetail.latitude", "articleDetail.longitude"], 'generate_map_image_placeholder', None),
}

# --- 데이터 가공 헬퍼 함수들 ---
def get_nested_value(data_dict, keys_string_or_list, default=""):
    if keys_string_or_list is None: return default
    if isinstance(keys_string_or_list, list):
        for keys_string in keys_string_or_list:
            val = get_single_nested_value(data_dict, keys_string, None)
            if val is not None and val != "": return val
        return default
    return get_single_nested_value(data_dict, keys_string_or_list, default)

def get_single_nested_value(data_dict, keys_string, default=""):
    if keys_string is None:
        return default
    keys = keys_string.split('.')
    value = data_dict
    try:
        for key_part in keys:
            if '[' in key_part and key_part.endswith(']'):
                key = key_part.split('[')[0]
                index = int(key_part.split('[')[1][:-1])
                if key in value and isinstance(value[key], list) and len(value[key]) > index:
                    value = value[key][index]
                else: return default
            elif isinstance(value, dict) and key_part in value:
                value = value[key_part]
            else: return default
        return value if value is not None else default
    except (KeyError, IndexError, TypeError, AttributeError): return default

def format_date_yyyy_mm(value, data_dict=None, json_path=None):
    if not value or len(str(value)) != 8: return ""
    try:
        dt_obj = datetime.datetime.strptime(str(value), "%Y%m%d")
        return f"{dt_obj.year}년 {dt_obj.month}월"
    except ValueError: return ""

def calculate_pyeong(m2_value, precision=1):
    if isinstance(m2_value, (int, float)) and m2_value > 0:
        return round(m2_value / 3.3058, precision)
    # return "" # This line is unreachable and can be removed if it exists.

def format_property_title_text(values, data_dict=None, json_path=None):
    # data_dict is expected to contain "매물순번" for dynamic numbering.
    prop_order = data_dict.get("매물순번", "N/A") if data_dict else "N/A"
    div_name = values[0] if len(values) > 0 and values[0] else ""
    apt_name = values[1] if len(values) > 1 and values[1] else ""
    return f"No.{prop_order} [{div_name}] {apt_name}"

def get_property_order(value=None, data_dict=None, json_path=None):
    # data_dict is expected to contain "매물순번" for dynamic page numbering.
    return data_dict.get("매물순번", "1") if data_dict else "1"

def format_household_count(value, data_dict=None, json_path=None):
    return f"{value}세대" if value else ""

def format_floor_count(value, data_dict=None, json_path=None):
    return f"{value}층" if value else ""

def format_heating_info(values, data_dict=None, json_path=None):
    method = values[0] if len(values) > 0 and values[0] else ""
    fuel = values[1] if len(values) > 1 and values[1] else ""
    if method and fuel: return f"{method}, {fuel}"
    return method or fuel or ""

def format_dong_ho_for_ppt(values, data_dict=None, json_path=None):
    building_name = values[0] if len(values) > 0 and values[0] else ""
    floor_info_raw = values[1] if len(values) > 1 and values[1] else ""
    
    floor_display = ""
    if isinstance(floor_info_raw, str):
        if "중층" in floor_info_raw or "고층" in floor_info_raw or "저층" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0]
        elif "/" in floor_info_raw:
            floor_display = floor_info_raw.split('/')[0] + "층"
        else: floor_display = floor_info_raw
    else: floor_display = str(floor_info_raw)
    
    return f"{building_name} {floor_display}".strip()

def format_supply_area_text(supply_m2_val, data_dict=None, json_path=None):
    supply_py_val = calculate_pyeong(supply_m2_val)
    return f"{supply_m2_val}㎡ ({supply_py_val}평)" if supply_m2_val else ""

def format_exclusive_area_text(exclusive_m2_val, data_dict=None, json_path=None):
    exclusive_py_val = calculate_pyeong(exclusive_m2_val)
    return f"{exclusive_m2_val}㎡ ({exclusive_py_val}평)" if exclusive_m2_val else ""

def format_room_bathroom_text(values, data_dict=None, json_path=None):
    room = values[0] if len(values) > 0 and values[0] else "0"
    bath = values[1] if len(values) > 1 and values[1] else "0"
    return f"{room} / {bath}"

def format_management_fee_for_ppt(value, data_dict=None, json_path=None):
    return "확인 어려움"

def format_price_details_for_ppt(values, data_dict=None, json_path=None):
    w_price_num = values[0] if len(values) > 0 and values[0] else 0
    r_price_num = values[1] if len(values) > 1 and values[1] else 0

    formatted_w = ""
    if isinstance(w_price_num, (int, float)) and w_price_num > 0:
        if w_price_num >= 10000:
            eok = int(w_price_num // 10000)
            manwon_part = int(w_price_num % 10000)
            formatted_w = f"{eok}억"
            if manwon_part > 0: formatted_w += f" {manwon_part:,}"
        else: formatted_w = f"{int(w_price_num):,}"
    
    formatted_r = ""
    if isinstance(r_price_num, (int, float)) and r_price_num > 0:
        formatted_r = f"{int(r_price_num):,}"
        
    if formatted_w and formatted_r: return f"{formatted_w} / {formatted_r} (만원)"
    if formatted_w: return f"{formatted_w} (만원)"
    if formatted_r: return f" / {formatted_r} (만원)"
    return "가격 정보 없음"

def get_first_available_value(values, data_dict=None, json_path=None):
    for val_candidate_key in values:
        val = val_candidate_key
        if val and val != "정보 없음": return str(val)
    return "내용 없음"

def format_note_for_ppt(values, data_dict=None, json_path=None):
    tags = values[0] if len(values) > 0 and isinstance(values[0], list) else []
    detail_desc = values[1] if len(values) > 1 and values[1] else ""
    manual_note = values[2] if len(values) > 2 and values[2] else ""

    if manual_note and manual_note != "정보 없음": return manual_note
    if tags: return ", ".join(tags)
    if detail_desc and detail_desc != "정보 없음": return detail_desc[:100] + ("..." if len(detail_desc) > 100 else "")
    return "특이사항 없음"

def prepare_image_url(value, data_dict=None, json_path=None):
    if value and isinstance(value, str): return IMAGE_BASE_URL + value
    return None

def prepare_image_url_if_type14(value, data_dict=None, json_path=None):
    grand_plan_list = get_nested_value(data_dict, 'articleDetail.grandPlanList', [])
    if grand_plan_list and isinstance(grand_plan_list, list) and len(grand_plan_list) > 0:
        first_plan = grand_plan_list[0]
        if isinstance(first_plan, dict) and first_plan.get('imageType') == "14" and first_plan.get('imageSrc'):
            return IMAGE_BASE_URL + first_plan.get('imageSrc')
    return None

def generate_map_image_placeholder(values, data_dict=None, json_path=None):
    lat = values[0] if len(values) > 0 and values[0] else None
    lon = values[1] if len(values) > 1 and values[1] else None
    if lat and lon: return f"지도 이미지 생성 필요 (위도: {lat}, 경도: {lon})"
    return "위경도 정보 없음"

helper_functions = {
    'format_date_yyyy_mm': format_date_yyyy_mm,
    'get_property_order': get_property_order,
    'format_property_title_text': format_property_title_text,
    'format_household_count': format_household_count,
    'format_floor_count': format_floor_count,
    'format_heating_info': format_heating_info,
    'format_dong_ho_for_ppt': format_dong_ho_for_ppt,
    'format_supply_area_text': format_supply_area_text,
    'format_exclusive_area_text': format_exclusive_area_text,
    'format_room_bathroom_text': format_room_bathroom_text,
    'format_management_fee_for_ppt': format_management_fee_for_ppt,
    'format_price_details_for_ppt': format_price_details_for_ppt,
    'get_first_available_value': get_first_available_value,
    'format_note_for_ppt': format_note_for_ppt,
    'prepare_image_url': prepare_image_url,
    'prepare_image_url_if_type14': prepare_image_url_if_type14,
    'generate_map_image_placeholder': generate_map_image_placeholder,
}

# --- PPTX 요소 분석 및 수정 함수 ---
def find_shape_by_name(slide_or_group, shape_name):
    for shape in slide_or_group.shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP: # MSO_SHAPE_TYPE.GROUP은 그룹 도형을 의미합니다.
            found_in_group = find_shape_by_name(shape, shape_name)
            if found_in_group:
                return found_in_group
    return None

def replace_text_in_shape(shape, placeholder_text_or_none, new_text, font_details=None):
    if not shape or not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    new_text_str = str(new_text if new_text is not None else "")

    text_frame.clear() 
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = new_text_str
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]

def add_image_to_shape(shape, image_url_or_path):
    if not shape: return
    # print(f"    '{shape.name}'에 이미지 채우기 시도: {image_url_or_path}")
    # 실제 이미지 로드 로직은 API 서버 환경에 맞게 재구성 필요 (예: 로컬 파일 접근, 외부 URL 접근 권한 등)
    # 현재는 주석 처리된 상태로 유지
    pass

def set_cell_text(cell, text, font_details=None):
    text_frame = cell.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    run = p.add_run()
    run.text = str(text if text is not None else "")
    
    if font_details:
        font = run.font
        if font_details.get("font_name"): font.name = font_details["font_name"]
        if font_details.get("font_size"): font.size = font_details["font_size"]
        if font_details.get("bold"): font.bold = font_details["bold"]

def fill_table_from_mappings(table_shape, table_mapping_rules, json_data, client_data):
    if not table_shape or not table_shape.has_table:
        # print(f"  경고: '{table_shape.name}'은(는) 표가 아니거나 찾을 수 없습니다.")
        return
    
    table = table_shape.table
    
    for (row_idx, col_idx), (header_text_or_none, json_path_or_list, func_name, font_details) in table_mapping_rules.items():
        if row_idx >= len(table.rows) or col_idx >= len(table.columns):
            # print(f"  경고: '{table_shape.name}' 표에 ({row_idx},{col_idx}) 셀이 없습니다.")
            continue
        
        cell = table.cell(row_idx, col_idx)
        final_text = ""
        
        if json_path_or_list == "TABLE_HEADER":
            final_text = header_text_or_none
        else:
            current_data_source = json_data
            actual_json_path = json_path_or_list
            
            raw_value = None
            if isinstance(actual_json_path, list):
                raw_values_list = []
                for jp in actual_json_path:
                    if jp in client_data:
                         raw_values_list.append(client_data.get(jp))
                    else:
                         raw_values_list.append(get_nested_value(json_data, jp))
                raw_value = raw_values_list
            elif actual_json_path :
                 if actual_json_path in client_data:
                     raw_value = client_data.get(actual_json_path)
                 else:
                     raw_value = get_nested_value(json_data, actual_json_path)

            final_text = raw_value
            if func_name and func_name in helper_functions:
                try:
                    final_text = helper_functions[func_name](raw_value, json_data, actual_json_path)
                except Exception as e:
                    # print(f"    오류: 함수 '{func_name}' 실행 중 오류 ({row_idx},{col_idx}): {e}")
                    final_text = "오류"
            
            if isinstance(final_text, list) and not func_name:
                final_text = ", ".join(map(str, filter(None, final_text)))
            elif final_text is None:
                final_text = ""
        
        set_cell_text(cell, str(final_text), font_details)
    # print(f"  '{table_shape.name}' 표 채우기 완료.")

# --- 슬라이드 복제 헬퍼 함수 ---
def clone_slide(prs, slide_index):
    """지정된 인덱스의 슬라이드를 복제하여 프레젠테이션에 추가합니다."""
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 모든 도형 복제
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # 테이블 복제
            old_table = shape.table
            new_table = new_slide.shapes.add_table(
                old_table.rows.count, 
                old_table.columns.count,
                shape.left, shape.top, shape.width, shape.height
            ).table
            
            # 셀 스타일 및 내용 복제
            for i, row in enumerate(old_table.rows):
                for j, cell in enumerate(row.cells):
                    new_cell = new_table.cell(i, j)
                    new_cell.fill.solid()
                    new_cell.fill.fore_color.rgb = cell.fill.fore_color.rgb
                    
                    # 텍스트 프레임의 속성(텍스트 박스) 복제
                    for paragraph in cell.text_frame.paragraphs:
                        new_paragraph = new_cell.text_frame.paragraphs[0] if i == 0 and j == 0 else new_cell.text_frame.add_paragraph()
                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run()
                            new_run.text = run.text
                            
                            # 서식 복제
                            if hasattr(run.font, 'bold'): new_run.font.bold = run.font.bold
                            if hasattr(run.font, 'italic'): new_run.font.italic = run.font.italic
                            if hasattr(run.font, 'underline'): new_run.font.underline = run.font.underline
                            if hasattr(run.font, 'size'): new_run.font.size = run.font.size
                            if hasattr(run.font, 'name'): new_run.font.name = run.font.name
                            if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'): 
                                new_run.font.color.rgb = run.font.color.rgb
            
            # 테이블 속성 복제
            for i in range(len(old_table.columns)):
                new_table.columns[i].width = old_table.columns[i].width
            for i in range(len(old_table.rows)):
                new_table.rows[i].height = old_table.rows[i].height
            
            # 테이블 이름 유지
            for shape_name in ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]:
                if shape.name == shape_name:
                    new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                    break
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 이미지 복제
            image_path = shape.image.filename if hasattr(shape.image, 'filename') else None
            if image_path:
                pic = new_slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                # 이미지 이름 유지
                for shape_name in ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:pic[last()]')[0].set('name', shape_name)
                        break
            else:
                # 이미지 경로를 찾을 수 없는 경우 빈 이미지 자리 유지
                placeholder = new_slide.shapes.add_shape(
                    MSO_SHAPE_TYPE.RECTANGLE, shape.left, shape.top, shape.width, shape.height
                )
                # 이름 유지
                for shape_name in ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                        break
        
        else:
            # 기타 도형 복제 (텍스트 박스 등)
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(left, top, width, height)
                new_text_frame = new_shape.text_frame
                old_text_frame = shape.text_frame
                
                # 텍스트 프레임 속성 복제
                new_text_frame.word_wrap = old_text_frame.word_wrap
                
                # 텍스트 내용 복제
                for i, paragraph in enumerate(old_text_frame.paragraphs):
                    new_paragraph = new_text_frame.paragraphs[0] if i == 0 else new_text_frame.add_paragraph()
                    new_paragraph.alignment = paragraph.alignment
                    
                    for run in paragraph.runs:
                        new_run = new_paragraph.add_run()
                        new_run.text = run.text
                        
                        # 서식 복제
                        if hasattr(run.font, 'bold'): new_run.font.bold = run.font.bold
                        if hasattr(run.font, 'italic'): new_run.font.italic = run.font.italic
                        if hasattr(run.font, 'size'): new_run.font.size = run.font.size
                        if hasattr(run.font, 'name'): new_run.font.name = run.font.name
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'): 
                            new_run.font.color.rgb = run.font.color.rgb
                
                # 텍스트박스 이름 유지
                for shape_name in ["txt_property_title", "txt_property_page"]:
                    if shape.name == shape_name:
                        new_slide.shapes._element.xpath('.//p:sp[last()]')[0].set('name', shape_name)
                        break
    
    # 슬라이드 속성 복제 (배경 등)
    if hasattr(source_slide, 'background') and hasattr(new_slide, 'background'):
        if hasattr(source_slide.background, 'fill') and hasattr(new_slide.background, 'fill'):
            if source_slide.background.fill.type == source_slide.background.fill.SOLID:
                new_slide.background.fill.solid()
                new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
    
    return new_slide

# --- 표지 슬라이드 채우기 함수 ---
def fill_cover_slide_data(slide, document_info, mappings_config):
    """
    표지 슬라이드의 내용을 채웁니다.
    Args:
        slide: 채울 슬라이드 객체
        document_info: 문서 정보 (client_data와 유사, "문서명", "고객명" 등 포함)
        mappings_config: 전체 매핑 규칙
    """
    print(f"[DEBUG] fill_cover_slide_data 호출됨. 슬라이드 ID: {slide.slide_id if slide else 'None'}, 문서명: {document_info.get('문서명', 'N/A')}")
    
    cover_slide_shape_keys = ["txt_document_title", "txt_client_name", "txt_company_name"]
    
    for shape_name in cover_slide_shape_keys:
        if shape_name not in mappings_config:
            print(f"[WARNING] 표지 매핑 '{shape_name}'이(가) 전체 매핑 규칙에 없습니다.")
            continue

        mapping_details = mappings_config[shape_name]
        if not (isinstance(mapping_details, tuple) and len(mapping_details) == 4):
            print(f"[WARNING] 표지 매핑 값 '{shape_name}'의 형식이 올바르지 않습니다.")
            continue
        
        placeholder_or_cell_coord, data_key, helper_func_name, font_details = mapping_details
        
        target_shape = find_shape_by_name(slide, shape_name)
        if not target_shape:
            print(f"[WARNING] 표지 슬라이드에서 '{shape_name}' 도형을 찾을 수 없습니다.")
            continue

        raw_value = document_info.get(data_key, f"키 '{data_key}' 없음")

        final_value = raw_value
        if helper_func_name and helper_func_name in helper_functions:
            try:
                final_value = helper_functions[helper_func_name](raw_value, document_info, data_key)
            except Exception as e:
                print(f"[ERROR] 표지 데이터에 헬퍼 함수 '{helper_func_name}' 적용 중 오류: {e}")
                final_value = "오류"
        
        if shape_name == "txt_client_name": # 특정 테이블 셀 처리 예시
            if target_shape.has_table and isinstance(placeholder_or_cell_coord, tuple) and len(placeholder_or_cell_coord) == 2:
                row_idx, col_idx = placeholder_or_cell_coord
                if row_idx < len(target_shape.table.rows) and col_idx < len(target_shape.table.columns):
                    set_cell_text(target_shape.table.cell(row_idx, col_idx), final_value, font_details)
                else:
                    print(f"[WARNING] '{shape_name}' 테이블에 ({row_idx},{col_idx}) 셀이 없습니다.")
            else:
                print(f"[WARNING] '{shape_name}' 도형은 테이블이 아니거나 셀 좌표가 부적절합니다.")
        elif target_shape.has_text_frame:
            replace_text_in_shape(target_shape, placeholder_or_cell_coord, final_value, font_details)
        else:
            print(f"[WARNING] '{shape_name}' 도형에 텍스트 프레임이 없습니다.")
    print(f"[DEBUG] fill_cover_slide_data 완료.")

# --- 매물 정보 슬라이드 채우기 함수 ---
def fill_property_data(slide, article_json_data, client_data, mappings_config):
    """
    단일 매물 정보 슬라이드의 내용을 채웁니다.
    Args:
        slide: 채울 슬라이드 객체
        article_json_data: 특정 매물에 대한 JSON 데이터 (반드시 "매물순번" 키 포함)
        client_data: 고객 및 문서 정보 (참고사항, 비고 등에 사용될 수 있음)
        mappings_config: 전체 매핑 규칙
    """
    property_order_str = str(article_json_data.get("매물순번", "N/A"))
    print(f"[DEBUG] fill_property_data 호출됨. 슬라이드 ID: {slide.slide_id if slide else 'None'}, 매물순번: {property_order_str}")

    # 제목 업데이트 (txt_property_title)
    title_shape_key = "txt_property_title"
    if title_shape_key in mappings_config:
        title_shape = find_shape_by_name(slide, title_shape_key)
        if title_shape:
            # format_property_title_text 헬퍼 함수를 사용하기 위해 article_json_data 전달
            _, data_keys, helper_name, font = mappings_config[title_shape_key]
            if helper_name == 'format_property_title_text' and isinstance(data_keys, list):
                values_for_helper = [get_nested_value(article_json_data, dk) for dk in data_keys]
                # 헬퍼 함수가 article_json_data 내의 "매물순번"을 참조하도록 수정했으므로, data_dict로 article_json_data 전달
                title_text = helper_functions[helper_name](values_for_helper, article_json_data, data_keys)
                if title_shape.has_table: # 제목이 테이블 셀 안에 있는 경우
                     coord = mappings_config[title_shape_key][0] # (row, col)
                     set_cell_text(title_shape.table.cell(coord[0], coord[1]), title_text, font)
                elif title_shape.has_text_frame: # 일반 텍스트 박스
                     replace_text_in_shape(title_shape, None, title_text, font)
        else:
            print(f"[WARNING] 매물 제목 도형 '{title_shape_key}'을(를) 슬라이드에서 찾을 수 없습니다.")

    # 페이지 번호 업데이트 (txt_property_page)
    page_shape_key = "txt_property_page"
    if page_shape_key in mappings_config:
        page_shape = find_shape_by_name(slide, page_shape_key)
        if page_shape and page_shape.has_text_frame:
            # get_property_order 헬퍼 함수가 article_json_data 내의 "매물순번"을 사용하도록 수정됨
            _, _, helper_name, font = mappings_config[page_shape_key]
            # page_text_val = helper_functions[helper_name](None, article_json_data, None)
            # replace_text_in_shape(page_shape, None, page_text_val, font)
            # 직접 매물 순번 사용 (get_property_order 헬퍼는 이제 article_json_data["매물순번"]을 반환)
            replace_text_in_shape(page_shape, None, property_order_str, font)
        else:
            print(f"[WARNING] 매물 페이지번호 도형 '{page_shape_key}'을(를) 찾을 수 없거나 텍스트 프레임이 없습니다.")

    # 테이블들 채우기
    table_keys = ["tbl_complex_info", "tbl_property_detail_1", "tbl_property_detail_2"]
    for table_key in table_keys:
        if table_key in mappings_config:
            table_shape = find_shape_by_name(slide, table_key)
            if table_shape and table_shape.has_table:
                fill_table_from_mappings(table_shape, mappings_config[table_key], article_json_data, client_data)
            else:
                print(f"[WARNING] 테이블 도형 '{table_key}'을(를) 찾을 수 없거나 테이블이 아닙니다.")
    
    # 이미지들 채우기
    image_keys = ["img_complex_view", "img_complex_floorplan", "img_complex_mapimg"]
    for img_key in image_keys:
        if img_key in mappings_config:
            img_shape = find_shape_by_name(slide, img_key)
            if img_shape:
                _, data_key_or_list, helper_name, _ = mappings_config[img_key]
                img_value_for_helper = get_nested_value(article_json_data, data_key_or_list)
                
                final_image_src_or_placeholder = img_value_for_helper # 기본값
                if helper_name and helper_name in helper_functions:
                    try:
                        # 헬퍼 함수는 (value, data_dict, json_path) 시그니처를 따름
                        final_image_src_or_placeholder = helper_functions[helper_name](img_value_for_helper, article_json_data, data_key_or_list)
                    except Exception as e_img_helper:
                        print(f"[ERROR] 이미지 헬퍼 함수 '{helper_name}' 실행 중 오류: {e_img_helper}")
                        final_image_src_or_placeholder = "이미지 오류"

                if final_image_src_or_placeholder:
                    if "지도 이미지 생성 필요" in str(final_image_src_or_placeholder) or \
                       "위경도 정보 없음" in str(final_image_src_or_placeholder) or \
                       "이미지 오류" in str(final_image_src_or_placeholder):
                        # 이미지 대신 플레이스홀더 텍스트 표시
                        if img_shape.has_text_frame:
                             replace_text_in_shape(img_shape, None, final_image_src_or_placeholder, {"font_name": "나눔고딕", "font_size": Pt(8)})
                        else: # 텍스트 프레임이 없는 경우 도형에 직접 추가 시도 (add_textbox_to_shape 같은 함수 필요)
                             print(f"[INFO] '{img_key}' 도형에 이미지 플레이스홀더 텍스트 '{final_image_src_or_placeholder}' 표시 시도 (텍스트 프레임 없음)")
                    else:
                        # 실제 이미지 로드 및 추가 로직 (현재 주석 처리)
                        # add_image_to_shape(img_shape, final_image_src_or_placeholder)
                        print(f"[INFO] 이미지 로드 대상: {img_key} <- {final_image_src_or_placeholder} (실제 로드는 주석 처리됨)")
                        # 임시로 이미지 경로를 텍스트로 표시 (디버깅용)
                        if img_shape.has_text_frame:
                            replace_text_in_shape(img_shape, None, str(final_image_src_or_placeholder), {"font_name": "나눔고딕", "font_size": Pt(8)})

            else:
                print(f"[WARNING] 이미지 도형 '{img_key}'을(를) 슬라이드에서 찾을 수 없습니다.")
    print(f"[DEBUG] fill_property_data 완료.")