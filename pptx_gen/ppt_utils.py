"""
PPT 조작 관련 유틸리티 함수들

이 모듈은 PPT 생성과 관련된 다양한 유틸리티 함수들을 포함합니다.
테스트 폴더의 슬라이드 복제 기능을 통합했습니다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.text.text import TextFrame
from typing import Optional, Tuple, Dict, Any
import re
import requests
from io import BytesIO
import six
from copy import deepcopy

def replace_text_in_frame(shape_or_cell_or_frame,
                         old_text_pattern: str,
                         new_text: str,
                         font_name: Optional[str] = None,
                         font_size: Optional[int] = None,
                         is_bold: Optional[bool] = None):
    """
    텍스트 프레임 내의 텍스트를 대체하고, 지정된 폰트 스타일을 적용합니다.
    old_text_pattern이 빈 문자열("")이면 전체 텍스트를 new_text로 교체하고 스타일을 적용합니다.
    """
    text_frame: Optional[TextFrame] = None
    if isinstance(shape_or_cell_or_frame, TextFrame):
        text_frame = shape_or_cell_or_frame
    elif hasattr(shape_or_cell_or_frame, 'text_frame') and shape_or_cell_or_frame.text_frame is not None:
        text_frame = shape_or_cell_or_frame.text_frame
    else:
        print(f"[ERROR replace_text_in_frame] Unexpected type or no text_frame for input object: {type(shape_or_cell_or_frame)}. Cannot get text_frame.")
        return

    if text_frame is None: # 위에서 할당되지 않은 경우를 대비한 최종 방어
        print(f"[ERROR replace_text_in_frame] text_frame is None for input object: {type(shape_or_cell_or_frame)}.")
        return

    # old_text_pattern이 실제 문자열일 때만 정규 표현식 컴파일
    old_text_pattern_re = None
    if old_text_pattern: # 빈 문자열이 아닐 경우에만 컴파일
        try:
            old_text_pattern_re = re.compile(re.escape(old_text_pattern))
        except re.error as e:
            print(f"[ERROR replace_text_in_frame] Invalid regex pattern from old_text_pattern '{old_text_pattern}': {e}")
            return # 패턴이 유효하지 않으면 더 이상 진행 불가

    for paragraph_idx, paragraph in enumerate(text_frame.paragraphs):
        current_new_text = str(new_text) if new_text is not None else ""
        text_changed = False

        if old_text_pattern == "" and old_text_pattern_re is None: # 플레이스홀더가 비어있으면, 단락 전체를 새 텍스트로 교체
            if paragraph.text != current_new_text:
                paragraph.text = current_new_text # 텍스트만 교체, run별 스타일링은 아래에서
                text_changed = True
            # 스타일은 항상 적용 (텍스트가 같더라도)
        elif old_text_pattern_re and old_text_pattern_re.search(paragraph.text): # 플레이스홀더가 있고, 해당 패턴이 단락에 존재
            original_paragraph_text = paragraph.text
            # 정규식 치환으로 새 텍스트 생성
            replaced_paragraph_text = old_text_pattern_re.sub(current_new_text, original_paragraph_text)

            if original_paragraph_text != replaced_paragraph_text:
                paragraph.text = replaced_paragraph_text # 텍스트만 교체
                text_changed = True
            # 스타일은 항상 적용 (텍스트가 같더라도, 패턴이 일치하는 한)
        elif not old_text_pattern and not paragraph.text and paragraph_idx == 0: # 플레이스홀더 없고, 단락도 비어있고, 첫번째 단락이면 (완전 빈칸에 텍스트 채우기)
            paragraph.text = current_new_text
            text_changed = True

        # 스타일 적용
        for run in paragraph.runs:
            if run.font:
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = Pt(font_size)
                if is_bold is not None:
                    run.font.bold = is_bold

        if text_changed:
            print(f"[DEBUG Style] Text SET and Applied Style (Font: {font_name or 'Default'}, Size: {font_size or 'Default'}pt, Bold: {is_bold if is_bold is not None else 'Default'}) for '{current_new_text[:30]}...'")
        elif old_text_pattern_re and old_text_pattern_re.search(paragraph.text): 
            print(f"[DEBUG Style] Style ONLY Applied (Font: {font_name or 'Default'}, Size: {font_size or 'Default'}pt, Bold: {is_bold if is_bold is not None else 'Default'}) for existing text matching '{old_text_pattern[:30]}...'")

def find_shape_by_placeholder_or_name(slide, identifier: str, is_placeholder_search: bool, slide_idx_for_log: int):
    """슬라이드에서 플레이스홀더 또는 이름으로 도형을 찾습니다."""
    if is_placeholder_search:
        # 플레이스홀더 텍스트가 정확히 일치하는 것을 먼저 찾고, 그 다음 포함된 것을 찾도록 개선 가능
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 정확히 일치하는 경우 (양끝 공백 제거 후 비교)
                if shape.text_frame.text and shape.text_frame.text.strip() == identifier:
                    print(f"[DEBUG][Slide {slide_idx_for_log+1}] Shape found by EXACT placeholder text '{identifier}': name='{shape.name if shape.name else 'N/A'}', type={shape.shape_type}")
                    return shape
                # 포함된 경우 (이전 로직)
                if identifier in shape.text_frame.text: # identifier가 비어있지 않아야 의미가 있음
                    print(f"[DEBUG][Slide {slide_idx_for_log+1}] Shape found by containing placeholder text '{identifier}': name='{shape.name if shape.name else 'N/A'}', type={shape.shape_type}")
                    return shape
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text_frame and cell.text_frame.text is not None:
                            if cell.text_frame.text.strip() == identifier:
                                print(f"[DEBUG find_shape] Found EXACT placeholder '{identifier}' in table cell ({row_idx},{col_idx}) on slide {slide_idx_for_log+1}")
                                return cell
                            if identifier in cell.text_frame.text:
                                print(f"[DEBUG find_shape] Found containing placeholder '{identifier}' in table cell ({row_idx},{col_idx}) on slide {slide_idx_for_log+1}")
                                return cell
    else: # 이름으로 검색
        try:
            shape_by_name = slide.shapes.get_by_name(identifier)
            print(f"[DEBUG][Slide {slide_idx_for_log+1}] Shape found by name '{identifier}': type={shape_by_name.shape_type}")
            return shape_by_name
        except KeyError: 
            pass
        except Exception as e: 
            print(f"[ERROR find_shape_by_placeholder_or_name] Error finding shape by name '{identifier}': {e}")

    print(f"[WARNING][Slide {slide_idx_for_log+1}] Shape with identifier '{identifier}' (search_mode: {'placeholder' if is_placeholder_search else 'name'}) not found. Skipping.")
    return None

def normalize_url(url: str, base_domain: str = "https://image.neonet.co.kr") -> str:
    """URL을 정규화합니다."""
    if not url: # 빈 URL 처리
        return ""
    if url.startswith("//"):
        return "https:" + url
    if url.startswith("/"):
        base = base_domain.rstrip('/')
        lpath = url.lstrip('/')
        return f"{base}/{lpath}"
    return url

def add_image_from_url(slide, image_url: str, left: float, top: float, width: Optional[float] = None, height: Optional[float] = None):
    """URL에서 이미지를 가져와 슬라이드에 추가합니다."""
    try:
        if not image_url:
            print("[WARNING] Empty image URL provided to add_image_from_url. Skipping.")
            return None

        normalized_url = normalize_url(image_url)
        if not normalized_url:
            print("[WARNING] Normalized image URL is empty. Skipping.")
            return None

        print(f"[DEBUG add_image_from_url] Downloading image from: {normalized_url}")

        response = requests.get(normalized_url, timeout=10)
        response.raise_for_status()

        image_stream = BytesIO(response.content)
        pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

        print(f"[DEBUG add_image_from_url] Image added successfully: position=({left}, {top}), size=({width}, {height}) if specified")
        return pic

    except requests.exceptions.RequestException as e:
        print(f"[ERROR add_image_from_url] Failed to download image from '{normalized_url}': {e}")
    except Exception as e:
        print(f"[ERROR add_image_from_url] Failed to add image for URL '{image_url}': {e}")
    return None

def apply_text_to_shape(shape, text_content: str, font_name: Optional[str] = None, font_size: Optional[int] = None, is_bold: Optional[bool] = None):
    """
    도형에 텍스트를 적용하고 지정된 폰트 스타일을 설정합니다.
    기존 텍스트는 모두 삭제됩니다.
    """
    text_frame: Optional[TextFrame] = None

    # TextFrame 가져오기 시도
    if isinstance(shape, TextFrame):
        text_frame = shape
    elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
        text_frame = shape.text_frame
    else:
        print(f"[ERROR apply_text_to_shape] Object {type(shape)} has no text_frame.")
        return False

    if text_frame is None:
        print("[ERROR apply_text_to_shape] Failed to get text_frame.")
        return False

    try:
        # 기존 텍스트 전체 대체
        text_frame.clear()  # 모든 단락 제거

        # 새 단락 추가 및 스타일링
        paragraph = text_frame.paragraphs[0]  
        paragraph.text = text_content

        # 폰트 스타일 설정
        for run in paragraph.runs:
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            if is_bold is not None:
                run.font.bold = is_bold

        return True
    except Exception as e:
        print(f"[ERROR apply_text_to_shape] Failed to apply text: {e}")
        return False

# --- 슬라이드 복제 기능 추가 ---
def duplicate_slide(prs, index):
    """슬라이드를 복제합니다. 이 함수는 널리 검증된 방법을 사용합니다.
    
    Args:
        prs: Presentation 객체
        index: 복제할 슬라이드 인덱스
    
    Returns:
        복제된 슬라이드 객체
    """
    if index >= len(prs.slides):
        print(f"[ERROR] 슬라이드 인덱스 {index}가 범위를 벗어났습니다. 총 슬라이드 수: {len(prs.slides)}")
        return None
    
    try:
        print(f"[INFO] 슬라이드 {index+1}번 복제 시작...")
        source = prs.slides[index]
        
        # 대상 슬라이드의 XML 표현을 가져옵니다
        slide_xml = source.element
        
        # 새 슬라이드 생성을 위한 XML 요소를 복제합니다
        new_slide_xml = deepcopy(slide_xml)
        
        # 복제된 슬라이드 XML을 프레젠테이션에 추가합니다
        # 복제된 슬라이드는 원본 후에 바로 추가됩니다
        # 상위 슬라이드 목록에서 삽입할 위치 찾기
        slide_collection = source.part.package.presentation_part.presentation.sldIdLst
        slide_id_list = [x for x in slide_collection]
        
        # 원본 슬라이드 다음에 새 슬라이드를 삽입
        new_slide_part = prs.part.package.presentation_part.new_slide_part(slide_xml)
        
        # 슬라이드 ID 생성
        rel_id = prs.part.package.presentation_part.relate_to(new_slide_part)
        
        # 새 슬라이드 ID 요소 생성 (ID 지정)
        new_slide_id = slide_collection._add_sldId(rel_id)
        
        # 원본 슬라이드 다음에 삽입
        for i, s in enumerate(slide_id_list):
            if s.rId == source.part.rId:
                slide_collection._insert_sldId(i+1, new_slide_id)
                break
        
        # 복제된 슬라이드 가져오기
        duplicate_slide = prs.slides[-1]  # 방금 추가된 슬라이드
        
        # 슬라이드 관계 복사
        for rel_id, rel in six.iteritems(source.part.rels):
            # 이미 복사된 관계 건너뛰기
            if rel.is_external or rel.reltype == duplicate_slide.part.reltype:
                continue
                
            # Target parts (이미지, 차트 등)를 새 슬라이드에 복사
            if hasattr(rel.target_part, 'partname'):
                duplicate_slide.part.relate_to(rel.target_part, rel.reltype)
        
        print(f"[INFO] 슬라이드 {index+1}번 복제 완료. 슬라이드 인덱스: {prs.slides.index(duplicate_slide)+1}")
        return duplicate_slide
    
    except Exception as e:
        import traceback
        print(f"[ERROR] 슬라이드 복제 중 오류가 발생했습니다: {e}")
        traceback.print_exc()
        return None
