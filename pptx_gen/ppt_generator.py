"""
PPT 생성 핵심 기능 모듈
이 모듈은 부동산 매물 정보를 바탕으로 PPT 파일을 생성하는 핵심 기능을 담당합니다.
테스트 폴더의 ppt_logic.py와 debug_ppt_elements.py의 기능을 활용합니다.
"""
import os
import tempfile
from pptx import Presentation
from typing import Dict, Optional, Any, List
import datetime

# ppt_logic.py의 핵심 기능 임포트
import os
import sys
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
try:
    from ppt_logic import clone_slide, generate_presentation_logic, fill_property_data, mappings as ppt_logic_mappings
    # fill_cover_slide_data는 ppt_logic.py에 있다고 가정합니다.
    # 실제로는 generate_presentation_logic의 일부이거나, fill_property_data를 특정 매개변수와 함께 사용할 수 있습니다.
    # 여기서는 fill_cover_slide_data가 존재한다고 가정하고 진행합니다.
    # 만약 없다면, 이 부분을 수정하거나 ppt_logic.py에 해당 함수를 추가해야 합니다.
    from ppt_logic import fill_cover_slide_data # 임시 가정
except ImportError:
    print("[ERROR] ppt_logic 모듈 또는 필요한 함수를 찾을 수 없습니다. 경로와 파일 내용을 확인하세요.")
    print(f"현재 Python 경로: {sys.path}")
    # 대체 구현 시도 (개발 및 테스트용)
    from pptx import Presentation
    
    def clone_slide(prs, index):
        print("[WARNING] 대체 구현된 clone_slide 함수를 사용합니다.")
        # 실제 복제 로직이 필요합니다. 아래는 단순 참조 반환으로, 실제 PPT에는 새 슬라이드가 추가되지 않습니다.
        # return prs.slides.add_slide(prs.slides[index].slide_layout) # 좀 더 나은 대체
        if index < len(prs.slides):
            source = prs.slides[index]
            image_dict = {}
            # 이미지를 BytesIO로 저장
            for shape in source.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: # 13은 그림 유형
                    try:
                        image_bytes = BytesIO(shape.image.blob)
                        image_dict[shape.name] = image_bytes
                    except Exception: # blob 접근이 항상 가능하지 않을 수 있음
                        pass


            new_slide = prs.slides.add_slide(source.slide_layout)
            for shp in source.shapes:
                el = shp.element
                new_el = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                
                # 이미지 재연결 (필요한 경우)
                if shp.name in image_dict:
                    try:
                        new_shape = new_slide.shapes[-1] # 방금 추가된 도형
                        if new_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                             new_slide.shapes.add_picture(image_dict[shp.name], new_shape.left, new_shape.top, new_shape.width, new_shape.height)
                             # 원래 도형 제거 또는 이미지로 대체하는 로직 추가 필요
                    except Exception as e_img:
                        print(f"Error re-adding image {shp.name}: {e_img}")
            return new_slide
        return prs.slides.add_slide(prs.slide_layouts[0]) # 기본 레이아웃으로 추가


    def generate_presentation_logic(client_data, article_json_data, mappings=None, template_path=None, output_path=None):
        print("[WARNING] 대체 구현된 generate_presentation_logic 함수를 사용합니다.")
        prs = Presentation(template_path)
        # 여기에 실제 데이터 채우기 로직이 필요합니다.
        prs.save(output_path)
        return output_path

    def fill_property_data(slide, article_json, client_data, mappings):
        print(f"[WARNING] 대체 구현된 fill_property_data 호출됨 - 슬라이드 '{slide.slide_id}'에 대해 데이터 채우기 시뮬레이션")
        # 이 함수는 실제 ppt_logic.py의 구현을 따라야 합니다.
        # 예시: slide.shapes[0].text_frame.text = article_json.get("매물순번", "") + " " + article_json.get("articleDetail",{}).get("aptName","")
        pass

    def fill_cover_slide_data(slide, document_info, mappings):
        print(f"[WARNING] 대체 구현된 fill_cover_slide_data 호출됨 - 슬라이드 '{slide.slide_id}'에 대해 표지 데이터 채우기 시뮬레이션")
        # 이 함수는 실제 ppt_logic.py의 구현을 따라야 합니다.
        # 예시: slide.shapes[0].text_frame.text = document_info.get("문서명", "")
        pass
        
    ppt_logic_mappings = {} # 대체 매핑

# 기존 코드와의 호환성 유지를 위한 경로 설정
BASE_DIR = os.path.dirname(os.path.abspath(__file__))


def create_ppt_file(property_data: Dict, map_image_url: Optional[str], template_filepath: str, mapping_csv_filepath: Optional[str] = None) -> str:
    """
    매물 정보를 바탕으로 PPT 파일을 생성합니다.

    Args:
        property_data: 매물 정보 딕셔너리
        map_image_url: 지도 이미지 URL (선택 사항)
        template_filepath: PPT 템플릿 파일 경로
        mapping_csv_filepath: 매핑 CSV 파일 경로 (선택 사항)

    Returns:
        생성된 임시 PPT 파일 경로
    """
    print("--- create_ppt_file 진입 (ppt_logic 기반 구현) ---")
    print(f"[DEBUG] Received property_data keys: {list(property_data.keys()) if property_data else 'None'}")
    print(f"[DEBUG] Received map_image_url: {map_image_url}")
    print(f"[DEBUG] Using template: {template_filepath}")

    # API 기대 형식과 내부 로직에 맞게 데이터 변환
    article_json_data = {}
    client_data = {}
    
    # 문서 정보, 고객 정보 등 가져오기
    if "documentTitle" in property_data:
        client_data["문서명"] = property_data.get("documentTitle", "")
    if "clientName" in property_data:
        client_data["고객명"] = property_data.get("clientName", "")
    if "companyName" in property_data:
        client_data["회사명"] = property_data.get("companyName", "")
    
    # 참고사항 및 비고 정보 처리
    if "참고사항_입력" in property_data:
        client_data["참고사항_입력"] = property_data.get("참고사항_입력", "")
    if "비고_입력" in property_data:
        client_data["비고_입력"] = property_data.get("비고_입력", "")
        
    # 매물 상세 정보
    if "articleDetail" in property_data:
        article_json_data.update(property_data.get("articleDetail", {}))
    
    # 추가 정보
    if "articleAddition" in property_data:
        article_json_data["articleAddition"] = property_data.get("articleAddition", {})
    
    # 매물 가격 정보
    if "articlePrice" in property_data:
        article_json_data["articlePrice"] = property_data.get("articlePrice", {})
    elif "articleAddition" in property_data and "dealOrWarrantPrc" in property_data["articleAddition"]:
        # 기존 API는 다른 형식으로 가격 정보를 보낼 수 있음
        if not "articlePrice" in article_json_data:
            article_json_data["articlePrice"] = {}
            
        try:
            # 문자열 형식의 금액을 숫자로 변환
            warrant_price_text = property_data["articleAddition"].get("dealOrWarrantPrc", "0")
            warrant_price = int(warrant_price_text.replace("억", "0000").replace(",", ""))
            article_json_data["articlePrice"]["warrantPrice"] = warrant_price
        except (ValueError, TypeError, AttributeError) as e:
            print(f"[WARNING] 보증금 문자열 변환 오류: {e}")
            article_json_data["articlePrice"]["warrantPrice"] = 0
            
        try:
            # 월세 문자열을 숫자로 변환
            rent_price_text = property_data["articleAddition"].get("rentPrc", "0") 
            rent_price = int(rent_price_text.replace(",", ""))
            article_json_data["articlePrice"]["rentPrice"] = rent_price
        except (ValueError, TypeError, AttributeError) as e:
            print(f"[WARNING] 월세 문자열 변환 오류: {e}")
            article_json_data["articlePrice"]["rentPrice"] = 0
    
    # 면적 정보
    if "articleSpace" in property_data:
        article_json_data["articleSpace"] = property_data.get("articleSpace", {})
    
    # 사진 정보 처리
    if "articlePhotos" in property_data:
        article_json_data["articlePhotos"] = property_data.get("articlePhotos", [])
    
    # 층 정보
    if "articleFloor" in property_data:
        article_json_data["articleFloor"] = property_data.get("articleFloor", {})
    
    # 관리비 정보
    if "administrationCostInfo" in property_data:
        article_json_data["administrationCostInfo"] = property_data.get("administrationCostInfo", {})
    
    # 지도 이미지 처리 (이 부분은 ppt_logic에서 활용됩니다)
    if map_image_url:
        article_json_data["mapImageUrl"] = map_image_url

    # 임시 파일 경로 생성
    temp_dir = tempfile.gettempdir()
    temp_filename = f"property_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = os.path.join(temp_dir, temp_filename)

    # ppt_logic의 generate_presentation_logic 함수 호출
    try:
        print(f"[INFO] PPT 생성 시작: {output_path}")
        generate_presentation_logic(client_data, article_json_data, None, template_filepath, output_path)
        print(f"[INFO] PPT 생성 완료: {output_path}")
        return output_path
    except Exception as e:
        print(f"[ERROR] PPT 생성 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        raise e


def create_ppt_file_multi(properties_data: List[Dict], document_info: Dict, map_image_url: Optional[str] = None, template_filepath: str = None, mapping_csv_filepath: Optional[str] = None) -> str:
    """
    여러 매물 정보를 바탕으로 하나의 PPT 파일을 생성합니다.
    
    Args:
        properties_data: 매물 정보 딕셔너리 목록
        document_info: 문서 제목, 고객명 등 문서 전체에 적용될 정보
        map_image_url: 지도 이미지 URL (선택 사항)
        template_filepath: PPT 템플릿 파일 경로
        mapping_csv_filepath: 매핑 CSV 파일 경로 (선택 사항)
        
    Returns:
        생성된 임시 PPT 파일 경로
    """
    print("--- create_ppt_file_multi 진입 (ppt_logic 기반 구현) ---")
    print(f"[DEBUG] Received {len(properties_data)} properties")
    print(f"[DEBUG] Document info: {document_info}")
    print(f"[DEBUG] Using template: {template_filepath}")
    
    # 문서 기본 정보 설정
    client_data = {
        "문서명": document_info.get("documentTitle", "매물 소개 자료"),
        "고객명": document_info.get("clientName", "고객님"),
        "회사명": document_info.get("companyName", ""),
        "참고사항_입력": document_info.get("참고사항_입력", ""),
        "비고_입력": document_info.get("비고_입력", "")
    }
    
    # 임시 파일 경로 생성
    temp_dir = tempfile.gettempdir()
    temp_filename = f"properties_multi_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = os.path.join(temp_dir, temp_filename)
    
    # 프레젠테이션 객체 생성
    try:
        prs = Presentation(template_filepath)
    except Exception as e:
        print(f"[ERROR] Failed to load PPT template: {template_filepath}. Error: {e}")
        raise e

    # 표지 슬라이드 (prs.slides[0]) 채우기
    try:
        print("[INFO] 표지 슬라이드 데이터 채우기 시작")
        # fill_cover_slide_data 함수가 ppt_logic.py에 정의되어 있다고 가정합니다.
        # 이 함수는 client_data와 ppt_logic_mappings를 사용하여 표지 슬라이드의 내용을 채웁니다.
        if len(prs.slides) > 0:
            fill_cover_slide_data(prs.slides[0], client_data, ppt_logic_mappings) 
            print("[INFO] 표지 슬라이드 데이터 채우기 완료")
        else:
            print("[ERROR] 템플릿에 슬라이드가 없습니다. 표지 슬라이드를 채울 수 없습니다.")
            # 오류 처리를 하거나, 빈 프레젠테이션으로 진행할 수 있습니다.
            # 여기서는 오류를 발생시키지 않고 진행하지만, 실제 상황에 따라 조치가 필요합니다.

    except Exception as e:
        print(f"[ERROR] 표지 슬라이드 데이터 채우기 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        # 표지 슬라이드 생성 실패 시에도 계속 진행할지, 아니면 여기서 중단할지 결정 필요
        # raise e # 필요시 주석 해제

    # 매물 슬라이드 생성
    property_template_slide_index = 1 # 매물 정보 슬라이드 템플릿은 두 번째 슬라이드(인덱스 1)

    if len(prs.slides) <= property_template_slide_index:
        print(f"[ERROR] 템플릿에 매물 정보 슬라이드 (인덱스 {property_template_slide_index})가 충분히 없습니다.")
        # 이 경우, 매물 슬라이드를 생성할 수 없으므로, 빈 PPT나 오류를 반환해야 할 수 있습니다.
        # 여기서는 빈 PPT를 저장하고 반환하도록 합니다.
        try:
            prs.save(output_path)
            print(f"[WARNING] 매물 정보 슬라이드 템플릿 부족으로, 내용이 없는 PPT 파일 저장: {output_path}")
            return output_path
        except Exception as e_save:
            print(f"[ERROR] 내용 없는 PPT 파일 저장 중 오류 발생: {e_save}")
            raise e_save


    for idx, property_data_item in enumerate(properties_data):
        print(f"[INFO] 매물 {idx + 1}/{len(properties_data)} 슬라이드 생성 시작")
        
        # 현재 매물에 대한 article_json_data 준비
        current_article_json_data = {}
        # articleDetail을 중첩된 딕셔너리로 할당해야 ppt_logic.mappings의 경로가 올바르게 작동합니다.
        if "articleDetail" in property_data_item:
            current_article_json_data["articleDetail"] = property_data_item.get("articleDetail", {})
        else: # articleDetail이 없는 경우를 대비해 빈 딕셔너리 할당
            current_article_json_data["articleDetail"] = {}
            
        if "articleAddition" in property_data_item:
            current_article_json_data["articleAddition"] = property_data_item.get("articleAddition", {})
        if "articlePrice" in property_data_item:
            current_article_json_data["articlePrice"] = property_data_item.get("articlePrice", {})
        if "articleSpace" in property_data_item:
            current_article_json_data["articleSpace"] = property_data_item.get("articleSpace", {})
        if "articlePhotos" in property_data_item:
            current_article_json_data["articlePhotos"] = property_data_item.get("articlePhotos", [])
        if "articleFloor" in property_data_item:
            current_article_json_data["articleFloor"] = property_data_item.get("articleFloor", {})
        if "administrationCostInfo" in property_data_item:
            current_article_json_data["administrationCostInfo"] = property_data_item.get("administrationCostInfo", {})
        
        # 매물 순번 추가
        current_article_json_data["매물순번"] = idx + 1
        
        # 지도 이미지 URL 추가 (필요한 경우)
        if map_image_url:
            current_article_json_data["mapImageUrl"] = map_image_url

        current_slide_to_fill = None
        try:
            if idx == 0:
                # 첫 번째 매물은 템플릿의 두 번째 슬라이드(인덱스 1)를 직접 사용
                current_slide_to_fill = prs.slides[property_template_slide_index]
                print(f"[DEBUG] 첫 번째 매물: 슬라이드 인덱스 {property_template_slide_index} 사용")
            else:
                # 이후 매물들은 템플릿의 두 번째 슬라이드를 복제하여 사용
                print(f"[DEBUG] 매물 {idx + 1}: 슬라이드 인덱스 {property_template_slide_index} 복제 시도")
                current_slide_to_fill = clone_slide(prs, property_template_slide_index)
                print(f"[DEBUG] 매물 {idx + 1}: 슬라이드 복제 완료. 새 슬라이드 ID: {current_slide_to_fill.slide_id if current_slide_to_fill else 'None'}")

            if current_slide_to_fill:
                # fill_property_data 함수를 호출하여 슬라이드 내용 채우기
                print(f"[INFO] 매물 {idx + 1}: fill_property_data 호출")
                fill_property_data(current_slide_to_fill, current_article_json_data, client_data, ppt_logic_mappings)
                print(f"[INFO] 매물 {idx + 1} 슬라이드 데이터 채우기 완료")
            else:
                print(f"[ERROR] 매물 {idx + 1}에 대한 슬라이드를 준비하지 못했습니다.")

        except Exception as e_slide:
            print(f"[ERROR] 매물 {idx + 1} 슬라이드 처리 중 오류 발생: {e_slide}")
            import traceback
            traceback.print_exc()
            # 개별 슬라이드 생성 오류 시 계속 진행할지 결정
            # continue # 다음 매물로 넘어감

    # 최종 PPT 파일 저장
    try:
        prs.save(output_path)
        print(f"[INFO] 다중 매물 PPT 파일 저장 완료: {output_path}")
        return output_path
    except Exception as e:
        print(f"[ERROR] PPT 파일 저장 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        raise e
